"""
PRANA CXO Deck v2 — Big4 consulting quality
Dense content, real data points, proper layout
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]

# ── Palette ──
N  = RGBColor(0x0F,0x0C,0x29)   # navy
IN = RGBColor(0x31,0x2E,0x81)   # indigo
VI = RGBColor(0x7C,0x3A,0xED)   # violet
PU = RGBColor(0x63,0x66,0xF1)   # purple
LA = RGBColor(0xA5,0xB4,0xFC)   # lavender
WH = RGBColor(0xFF,0xFF,0xFF)
SL = RGBColor(0x0F,0x17,0x2A)   # slate-900
G5 = RGBColor(0x64,0x74,0x8B)   # slate-500
G2 = RGBColor(0xF8,0xFA,0xFF)   # near-white
GR = RGBColor(0x05,0x96,0x69)   # green
LG = RGBColor(0xEC,0xFD,0xF5)   # light-green
AM = RGBColor(0xF5,0x9E,0x0B)   # amber
RE = RGBColor(0xE1,0x1D,0x48)   # red
CY = RGBColor(0x08,0x91,0xB2)   # cyan
IL = RGBColor(0xEE,0xF2,0xFF)   # indigo-light
DK = RGBColor(0x1E,0x1B,0x4B)   # dark-indigo
BG = RGBColor(0xF1,0xF5,0xF9)   # slate-100
R2 = RGBColor(0xFF,0xF1,0xF2)
A2 = RGBColor(0xFF,0xFB,0xEB)
V2 = RGBColor(0xF5,0xF3,0xFF)
C2 = RGBColor(0xF0,0xF9,0xFF)
P2 = RGBColor(0x4F,0x46,0xE5)   # indigo-600

# ── Primitives ──
def R(slide, l, t, w, h, fill=None):
    s = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    s.line.fill.background()
    if fill: s.fill.solid(); s.fill.fore_color.rgb = fill
    else:    s.fill.background()
    return s

def T(slide, text, l, t, w, h, sz=9.5, bold=False, color=SL,
      align=PP_ALIGN.LEFT, italic=False):
    txb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    txb.word_wrap = True
    tf = txb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run()
    r.text = text; r.font.size = Pt(sz)
    r.font.bold = bold; r.font.color.rgb = color; r.font.italic = italic
    return txb

def TT(slide, lines, l, t, w, h, sz=9.5, color=SL, leading=0.22):
    """Multi-line text with tight leading"""
    txb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    txb.word_wrap = True
    tf = txb.text_frame; tf.word_wrap = True
    for i, (bold, txt) in enumerate(lines):
        p = tf.paragraphs[i] if i==0 else tf.add_paragraph()
        r = p.add_run(); r.text = txt
        r.font.size = Pt(sz); r.font.bold = bold; r.font.color.rgb = color
    return txb

# ── Logo: coloured P-box + RANA text ──
def logo(slide, x=0.22, y=0.13):
    R(slide, x, y, 0.28, 0.28, fill=PU)
    T(slide, "P", x+0.055, y+0.02, 0.18, 0.24, sz=11, bold=True, color=WH, align=PP_ALIGN.CENTER)
    T(slide, "RANA", x+0.32, y+0.04, 0.7, 0.24, sz=10, bold=True, color=IN)

# ── Slide header: logo + divider line + title (same row) ──
def header(slide, title, badge=""):
    R(slide, 0, 0, 13.33, 0.55, fill=WH)
    R(slide, 0, 0.55, 13.33, 0.025, fill=IL)
    logo(slide, 0.22, 0.13)
    if badge:
        T(slide, badge, 1.05, 0.15, 2.5, 0.2, sz=6.5, bold=True, color=PU)
        T(slide, title, 1.05, 0.27, 11.8, 0.26, sz=14, bold=True, color=SL)
    else:
        T(slide, title, 1.05, 0.13, 11.8, 0.38, sz=14, bold=True, color=SL)

# ── Footer ──
def footer(slide, n, total=15):
    R(slide, 0, 7.2, 13.33, 0.3, fill=BG)
    R(slide, 0, 7.195, 13.33, 0.008, fill=IL)
    T(slide, "Private & Confidential", 0.22, 7.22, 4, 0.22, sz=7, color=G5)
    T(slide, "NS", 6.3, 7.22, 0.8, 0.22, sz=7, color=G5, align=PP_ALIGN.CENTER)
    T(slide, f"Slide {n} of {total}", 9.3, 7.22, 3.8, 0.22, sz=7, color=G5, align=PP_ALIGN.RIGHT)

# ── Body area y-start after header ──
BY = 0.65   # body starts at 0.65"
BH = 6.45   # body height to footer

def bg(slide):
    R(slide, 0, 0, 13.33, 7.5, fill=G2)

# ── Stat block: big number + label + delta ──
def stat(slide, x, y, w, h, val, label, delta="", dcol=GR, bg_col=WH):
    R(slide, x, y, w, h, fill=bg_col)
    T(slide, val,   x+0.12, y+0.1,  w-0.2, 0.42, sz=22, bold=True, color=SL)
    T(slide, label, x+0.12, y+0.52, w-0.2, 0.2,  sz=8,  color=G5)
    if delta:
        T(slide, delta, x+0.12, y+0.72, w-0.2, 0.2, sz=7.5, color=dcol)

# ── Thin metric row inside a box ──
def mrow(slide, x, y, w, label, value, col=PU):
    T(slide, label, x+0.12, y, w*0.7,  0.22, sz=8.5, color=G5)
    T(slide, value, x+w*0.72, y, w*0.28, 0.22, sz=9,  bold=True, color=col, align=PP_ALIGN.RIGHT)

# ── Gradient-style section box ──
def sbox(slide, x, y, w, h, title, body, bg_col=IL, tc=IN, bc=SL):
    R(slide, x, y, w, h, fill=bg_col)
    R(slide, x, y, 0.03, h, fill=PU)   # left accent bar
    T(slide, title, x+0.12, y+0.08, w-0.2, 0.24, sz=9.5, bold=True, color=tc)
    T(slide, body,  x+0.12, y+0.35, w-0.18, h-0.45, sz=8.5, color=bc)

# ── Bullet row with coloured dot ──
def brow(slide, x, y, w, text, dot=PU, sz=8.5):
    R(slide, x, y+0.065, 0.07, 0.07, fill=dot)
    T(slide, text, x+0.14, y, w-0.14, 0.26, sz=sz, color=SL)

# ── Bar row ──
def bar(slide, x, y, bw, label, pct, value_lbl, col=PU):
    T(slide, label, x, y, 2.0, 0.22, sz=8.5, color=G5)
    R(slide, x+2.08, y+0.05, bw, 0.15, fill=BG)
    if pct > 0:
        R(slide, x+2.08, y+0.05, bw*(pct/100), 0.15, fill=col)
    T(slide, value_lbl, x+2.08+bw+0.08, y, 0.7, 0.22, sz=8.5, bold=True, color=col, align=PP_ALIGN.RIGHT)

# ═══════════════════════════════════════════
# SLIDE 1 — COVER
# ═══════════════════════════════════════════
s1 = prs.slides.add_slide(BLANK)
R(s1, 0, 0, 13.33, 7.5, fill=N)
R(s1, 0, 7.28, 13.33, 0.22, fill=PU)  # bottom stripe

# Sanskrit strip — top left
R(s1, 0.4, 0.32, 9.5, 0.52, fill=RGBColor(0x2D,0x2B,0x5E))
T(s1, "प्राण", 0.52, 0.35, 1.1, 0.42, sz=14, color=LA)
R(s1, 1.72, 0.42, 0.015, 0.28, fill=RGBColor(0x4B,0x48,0x7A))
T(s1, "Prāṇa — Sanskrit for life force, the vital energy that sustains all existence. "
       "Just as प्राण is essential to life — PRANA is essential to our career.",
  1.78, 0.38, 8.0, 0.42, sz=8.5, color=RGBColor(0xC4,0xB5,0xFD))

# Hero headline
T(s1, "Every worker deserves proof of", 0.4, 1.05, 11, 0.32, sz=12, color=RGBColor(0xC4,0xB5,0xFD))
T(s1, "their career.", 0.4, 1.35, 10, 0.82, sz=44, bold=True, color=WH)
R(s1, 0.4, 2.22, 0.65, 0.05, fill=PU)

# Three value bullets
bullets_c = [
    "Employer-Pushed  →  AI-Processed  →  Employee-Owned  career document vault",
    "Global Alumni Portal — your vault survives your resignation. Always on. Always yours.",
    "Built for   Trust  ·  Compliance  ·  Permanence",
]
for i, b in enumerate(bullets_c):
    R(s1, 0.4, 2.34+i*0.48, 0.06, 0.06, fill=PU)
    T(s1, b, 0.56, 2.30+i*0.48, 10.5, 0.38, sz=10.5, color=LA if i!=1 else WH)

# PRANA full form — right column
ff = [("P","ersonal"),("R","epository and"),("A","uthenticated"),("N","etworked"),("A","uthentications")]
for i,(letter,rest) in enumerate(ff):
    y = 3.82+i*0.44
    R(s1, 0.4, y+0.03, 0.28, 0.28, fill=RGBColor(0x2D,0x2B,0x5E))
    T(s1, letter, 0.48, y+0.04, 0.2, 0.24, sz=11, bold=True, color=LA, align=PP_ALIGN.CENTER)
    T(s1, rest, 0.74, y+0.05, 3.0, 0.24, sz=9, color=RGBColor(0x94,0xA3,0xB8))

# PRANA logo top right
R(s1, 12.4, 0.18, 0.3, 0.3, fill=PU)
T(s1, "P", 12.44, 0.2, 0.22, 0.26, sz=11, bold=True, color=WH, align=PP_ALIGN.CENTER)
T(s1, "RANA", 12.75, 0.22, 0.72, 0.26, sz=10, bold=True, color=WH)

# Prepared by
T(s1, "Prepared by: Nilesh Shah    ·    Classification: Private & Confidential",
  0.4, 7.0, 10, 0.22, sz=8, color=RGBColor(0x64,0x74,0x8B))

# ═══════════════════════════════════════════
# SLIDE 2 — TABLE OF CONTENTS
# ═══════════════════════════════════════════
s2 = prs.slides.add_slide(BLANK)
bg(s2)
header(s2, "What we'll cover today")
footer(s2, 2)

R(s2, 0.22, BY, 12.88, 0.42, fill=DK)
T(s2, "15 slides  ·  ~20 minutes  ·  Prepared for CXO audience", 0.42, BY+0.1, 10, 0.26, sz=9, color=LA)

toc_items = [
    ("03","The Problem","₹42,000 Cr annual loss · 15 crore workers with no trail"),
    ("04","Platform Architecture","Employer → AI Engine → Employee · 202 async · dual-region"),
    ("05","The Seven Trust Pillars","Zero-knowledge · per-DEK · immutable audit · post-exit fence"),
    ("06","Alumni Portal","Career vault that survives resignation · 38-sec retrieval"),
    ("07","CHRO Dashboard","91.6% compliance · 0 manual digest steps · 1-click audit export"),
    ("08","CFO Dashboard","99.4% TDS alignment · ₹50Cr penalty structurally impossible"),
    ("09","CISO Dashboard","0 plaintext PAN in DB · 2ms destruction · RLS-enforced"),
    ("10","For the Employee","6 DPDP rights built-in · ₹0 bank verification · 0 employer post-exit access"),
    ("11","AI Pipeline + Guardrails","6 stages · local LLM · privacy filter · multilingual"),
    ("12","Labour Law Compliance","4 Codes · 29 Acts · Gratuity Form F/I/L/M · Maternity"),
    ("13","DPDP Act 2023","All 6 rights · per-purpose consent · ₹250 Cr penalty avoided"),
    ("14","Security + Scale","7 trust layers · 1 crore employees · 125 PB · 99.99% uptime"),
    ("15","Immutable Audit + Trust","Append-only · 7-year retention · watermarked shares"),
    ("16","Market Landscape","DigiLocker gap · HRMS blind spot · where PRANA fits"),
    ("17","Why Now + Next Step","3 converging forces · 30-day pilot offer"),
]
# 2 cols
col_x = [0.22, 6.9]
rows_per_col = 8
for i,(num,title,sub) in enumerate(toc_items):
    col = 0 if i < rows_per_col else 1
    row = i if i < rows_per_col else i - rows_per_col
    x = col_x[col]
    y = BY+0.58+row*0.72
    if y > 7.1: break
    R(s2, x, y, 0.3, 0.3, fill=PU)
    T(s2, num, x+0.04, y+0.04, 0.26, 0.24, sz=8.5, bold=True, color=WH, align=PP_ALIGN.CENTER)
    T(s2, title, x+0.4, y, 6.0, 0.24, sz=10, bold=True, color=SL)
    T(s2, sub, x+0.4, y+0.26, 6.0, 0.22, sz=7.5, color=G5)

# ═══════════════════════════════════════════
# SLIDE 3 — THE PROBLEM
# ═══════════════════════════════════════════
s3 = prs.slides.add_slide(BLANK)
bg(s3)
header(s3, "India's workforce runs on documents no one can find.", "THE PROBLEM")
footer(s3, 3)

# 3 persona cards
personas = [
    ("The Worker","No Portability",
     ['"My old employer won\'t release my salary slips. 3 years of career evidence — gone."',
      "15 crore gig & contract workers change jobs with zero portable document trail.",
      "Average time to recover 3 years of salary slips from 2 ex-employers: 6–8 weeks.",
      "Banks reject loan applications. Visa officers reject applications. All for want of a PDF."],
     R2, RE),
    ("The CHRO","No Visibility",
     ['"We run SAP, Darwinbox, and 2 local HRMS. Zero single view of document coverage."',
      "Audit season means manual spreadsheets, helpdesk tickets, and executive escalations.",
      "87% vault completeness sounds good — until you find out the 13% is all Finance dept.",
      "Compliance calendar is a shared Excel file someone updates once a quarter."],
     A2, AM),
    ("The CISO / InfoSec","No Trail",
     ['"Who accessed which salary slip, from which IP, at 2 AM? We have zero answer."',
      "No watermark on documents leaves no exfiltration evidence.",
      "No per-employee encryption means one breach = entire workforce data exposed.",
      "DPDP Act 2023 fines up to ₹250 Cr. Audit trail: a folder of email threads."],
     IL, PU),
]
for i,(role,badge,pts,bg_c,ac) in enumerate(personas):
    x = 0.22 + i*4.37
    R(s3, x, BY, 4.2, 5.05, fill=bg_c)
    R(s3, x, BY, 4.2, 0.04, fill=ac)
    R(s3, x, BY, 0.03, 5.05, fill=ac)
    T(s3, role, x+0.12, BY+0.1, 4.0, 0.28, sz=11, bold=True, color=ac)
    R(s3, x+0.12, BY+0.42, 1.4, 0.22, fill=ac)
    T(s3, badge, x+0.12, BY+0.44, 1.4, 0.2, sz=7.5, bold=True, color=WH, align=PP_ALIGN.CENTER)
    for j,pt in enumerate(pts):
        R(s3, x+0.12, BY+0.82+j*0.98+0.04, 0.06, 0.06, fill=ac)
        T(s3, pt, x+0.25, BY+0.82+j*0.98, 3.75, 0.85, sz=8.5,
          color=SL, italic=(j==0))

# Bottom bar
R(s3, 0.22, BY+5.2, 12.88, 0.88, fill=RE)
T(s3, "₹42,000 Cr", 0.42, BY+5.25, 3.2, 0.42, sz=26, bold=True, color=WH)
T(s3, "estimated annual enterprise loss from document fraud,\nre-issuance delays, and non-compliance penalties",
  3.7, BY+5.3, 5.5, 0.6, sz=9, color=RGBColor(0xFF,0xC5,0xC5))
R(s3, 9.3, BY+5.25, 0.015, 0.65, fill=RGBColor(0xFF,0x80,0x80))
T(s3, "15 Crore", 9.45, BY+5.25, 1.8, 0.38, sz=22, bold=True, color=WH)
T(s3, "gig & contract workers with\nzero portable document trail", 9.45, BY+5.62, 3.4, 0.38, sz=8.5, color=RGBColor(0xFF,0xC5,0xC5))

# ═══════════════════════════════════════════
# SLIDE 4 — PLATFORM ARCHITECTURE
# ═══════════════════════════════════════════
s4 = prs.slides.add_slide(BLANK)
bg(s4)
header(s4, "Employer pushes. AI processes. Employee owns — permanently.", "PLATFORM ARCHITECTURE")
footer(s4, 4)

# Three zones
zones = [
    ("EMPLOYER", "Portal Upload\nHRMS API Push (SAP, Darwinbox,\nKeka, Zoho, Excel)\nBatch ZIP + CSV manifest\n\nWebhook confirmation\nIdemptotent — no duplicate risk",
     IN, IL, 0.22),
    ("PRANA VAULT ENGINE", "Kafka · Temporal · YugabyteDB (dual-region)\nAWS KMS · Redis CRDT · Tesseract / Textract\n\nHTTP contract:\nValidate → S3 put → 1 DB write\n→ 1 Kafka publish → 202 ack\n\nAll pipeline work is asynchronous.\nHTTP handler returns in milliseconds.\nFive Kafka consumers own fan-out:\nAudit · Workflow · SSE · Notify · Analytics",
     WH, DK, 4.5),
    ("EMPLOYEE", "Mobile App (iOS + Android)\nAI-powered career insights\nWatermarked secure share\nAlumni access — post-exit\nDPDP rights self-service\nC-Share: named, expiring, revocable",
     GR, LG, 10.0),
]
for i,(title,body,tc,bg_c,x) in enumerate(zones):
    w = 4.15 if i==1 else 3.3
    h = 4.5
    R(s4, x, BY, w, h, fill=bg_c)
    R(s4, x, BY, w, 0.32, fill=tc if i!=1 else RGBColor(0x4F,0x46,0xE5))
    T(s4, title, x+0.12, BY+0.06, w-0.2, 0.24, sz=9, bold=True,
      color=WH if i in [0,2] else LA,align=PP_ALIGN.CENTER if i==1 else PP_ALIGN.LEFT)
    T(s4, body, x+0.12, BY+0.42, w-0.22, h-0.55, sz=8.5,
      color=WH if i==1 else (WH if i==2 else SL))
    if i < 2:
        T(s4, "→", x+w+0.08, BY+2.0, 0.4, 0.4, sz=24, color=G5, align=PP_ALIGN.CENTER)

# Technical guarantees strip
R(s4, 0.22, BY+4.65, 12.88, 1.7, fill=WH)
R(s4, 0.22, BY+4.65, 12.88, 0.025, fill=IL)
guarantees = [
    ("⚡ 202 in <50ms","HTTP handler exits immediately. No synchronous pipeline blocking. HRMS never waits."),
    ("🔑 Envelope Encryption","Employee DEK → Tenant KEK → AWS KMS ap-south-1. Zero plaintext in DB or cache."),
    ("🌐 Dual-Region","YugabyteDB: Mumbai + Hyderabad. Active-active. Sub-10ms sync. 99.99% SLA."),
    ("🔄 53 Temporal Workflows","Durable, retryable, auditable. Zero cron jobs. Zero Celery. Zero polling."),
]
for i,(t,b) in enumerate(guarantees):
    x = 0.35 + i*3.25
    T(s4, t, x, BY+4.78, 3.1, 0.24, sz=9, bold=True, color=IN)
    T(s4, b, x, BY+5.06, 3.1, 1.1, sz=8, color=G5)

# ═══════════════════════════════════════════
# SLIDE 5 — SEVEN TRUST PILLARS
# ═══════════════════════════════════════════
s5 = prs.slides.add_slide(BLANK)
bg(s5)
header(s5, "Seven structural guarantees — not policy promises, but architecture.", "TRUST ARCHITECTURE")
footer(s5, 5)

pillars = [
    ("01","Vault Ownership Transfers on Routing","STRUCTURAL",
     "Once routed, employer CANNOT open, download, share, or delete. Not a permission flag — the database row changes owner."),
    ("02","One Person. One Key. One Vault.","CRYPTOGRAPHIC",
     "Each employee has their own DEK. One key compromise = one person's blast radius. Full DB + no KMS = zero readable data."),
    ("03","The System Doesn't Know Who You Are","ZERO-KNOWLEDGE",
     "PAN → HMAC-SHA256 in 2ms. Plaintext destroyed immediately. No PAN in any DB row, cache entry, or log line. Ever."),
    ("04","Sharing Requires Your Action. Always.","CONSENT-GATED",
     "Only the employee creates share tokens. Named recipient. Time-limited. Revocable. Watermarked. View-count visible to owner."),
    ("05","Portal Admin Cannot Read Documents","RLS ENFORCED",
     "Zero SELECT on document rows. PostgreSQL RLS refuses the query at database level — not application logic, database law."),
    ("06","Every Access Visible to Document Owner","TRANSPARENT",
     "Actor, timestamp, IP, channel, watermark status — all in employee activity log. No black boxes. No hidden access."),
    ("07","Your Vault Follows You, Not Your Employer","PORTABLE",
     "On resignation, employer access revoked immediately. Vault URL is permanent. Alumni access in 38 seconds."),
]
cols, rows = 4, 2
for i, (num,title,badge,body) in enumerate(pillars):
    if i < 4:
        x = 0.22 + i*3.28; y = BY
    else:
        x = 0.22 + (i-4)*4.37; y = BY+2.9
    w = 3.1 if i < 4 else 4.1
    h = 2.6
    R(s5, x, y, w, h, fill=WH)
    R(s5, x, y, w, 0.025, fill=PU)
    R(s5, x, y, 0.025, h, fill=PU)
    R(s5, x+0.1, y+0.1, 0.26, 0.26, fill=DK)
    T(s5, num, x+0.1, y+0.1, 0.26, 0.26, sz=9, bold=True, color=LA, align=PP_ALIGN.CENTER)
    R(s5, x+0.44, y+0.14, 1.4, 0.2, fill=IL)
    T(s5, badge, x+0.44, y+0.15, 1.4, 0.2, sz=6.5, bold=True, color=PU, align=PP_ALIGN.CENTER)
    T(s5, title, x+0.1, y+0.42, w-0.18, 0.38, sz=9.5, bold=True, color=SL)
    T(s5, body, x+0.1, y+0.85, w-0.18, 1.65, sz=8.5, color=G5)

# ═══════════════════════════════════════════
# SLIDE 6 — ALUMNI PORTAL
# ═══════════════════════════════════════════
s6 = prs.slides.add_slide(BLANK)
bg(s6)
header(s6, "Your career vault survives your resignation. Always on. Always yours.", "ALUMNI PORTAL")
footer(s6, 6)

# Left: feature list
features6 = [
    ("Post-Exit Access in 38 Seconds",
     "Ex-employee authenticates with phone OTP. Salary slips, Form 16, offer letters, relieving letter across every employer — in one place. No HR call. No email. No waiting."),
    ("Portable Across Every Employer",
     "Career spanning Vertex → Indigo → Bluestar? One vault. All employers' documents co-exist, tenant-isolated, employee-unified. 3 employers, 1 login."),
    ("Employer Access Revoked Immediately on Exit",
     "On resignation flag, employer loses all access instantly. Alumni vault documents are classified employee property. Employer query returns ACCESS_DENIED — at the database row level."),
    ("Self-Serve Document Sharing",
     "Bank loan needs salary slip proof? Employee creates a QR-watermarked share link. Bank scans QR — verified. No phone call to HR. No 3-day wait. ₹0 cost."),
    ("Nomination — First Platform in India",
     "Employee nominates family member for vault access. Built-in under DPDP S.14. In the event of incapacitation, nominee accesses with verified credentials."),
]
for i,(t,b) in enumerate(features6):
    y = BY + i*1.1
    R(s6, 0.22, y, 7.5, 1.0, fill=WH)
    R(s6, 0.22, y, 0.03, 1.0, fill=PU)
    T(s6, t, 0.38, y+0.07, 7.1, 0.26, sz=9.5, bold=True, color=IN)
    T(s6, b, 0.38, y+0.38, 7.1, 0.55, sz=8.5, color=G5)

# Right: stat panel
R(s6, 8.02, BY, 5.08, 5.5, fill=DK)
R(s6, 8.02, BY, 5.08, 0.025, fill=PU)
metrics6 = [
    ("38 sec","Avg alumni document retrieval — any employer"),
    ("0","HR calls needed for alumni access — ever"),
    ("90+","System-generated notifications across 6 channels"),
    ("₹0","Cost to employee for bank document verification"),
    ("2ms","PAN lifetime in memory — then HMAC, then destroyed"),
    ("0","Employer accesses post-exit — structurally impossible"),
    ("3","Employers, 1 vault, 1 login — career is portable"),
]
T(s6, "By the numbers", 8.22, BY+0.1, 4.6, 0.28, sz=9, bold=True, color=LA)
for i,(v,l) in enumerate(metrics6):
    y = BY+0.52+i*0.7
    R(s6, 8.22, y, 4.6, 0.62, fill=RGBColor(0x2D,0x2B,0x5E))
    T(s6, v, 8.35, y+0.05, 1.5, 0.42, sz=18, bold=True, color=LA)
    T(s6, l, 9.95, y+0.1, 2.7, 0.38, sz=8, color=RGBColor(0x94,0xA3,0xB8))

# Notif channels pill row
R(s6, 8.02, BY+5.25, 5.08, 0.5, fill=IN)
channels = ["Email","SMS","WhatsApp","Push","Portal Bell","Incident Log"]
T(s6, "  ·  ".join(channels), 8.22, BY+5.35, 4.7, 0.3, sz=8, color=LA)

# ═══════════════════════════════════════════
# SLIDE 7 — CHRO
# ═══════════════════════════════════════════
s7 = prs.slides.add_slide(BLANK)
bg(s7)
header(s7, "Every HR obligation tracked. Every audit instant. Zero manual steps.", "FOR THE CHRO")
footer(s7, 7)

# Left: digest mock
R(s7, 0.22, BY, 7.7, 6.25, fill=WH)
R(s7, 0.22, BY, 7.7, 0.025, fill=PU)
T(s7, "CHRO Weekly Vault Snapshot  ·  Week 32  ·  Aug 2025", 0.38, BY+0.1, 7.3, 0.28, sz=10, bold=True, color=SL)
T(s7, "TechServe India  ·  2,000 employees  ·  PRANA live 7 months", 0.38, BY+0.4, 7.3, 0.2, sz=8, color=G5)

# Stat row
stats7 = [("214","Docs pushed","+ 18 vs last week",GR),("87%","Vault health","↑ 1.2% WoW",GR),("91.6%","Compliance","vs 89.1% target",GR),("4","Exceptions","2 past 4-hr SLA",RE)]
for i,(v,l,d,dc) in enumerate(stats7):
    x = 0.35+i*1.88
    R(s7, x, BY+0.66, 1.76, 1.0, fill=BG)
    T(s7, v, x+0.1, BY+0.72, 1.6, 0.42, sz=20, bold=True, color=SL)
    T(s7, l, x+0.1, BY+1.14, 1.6, 0.2, sz=7.5, color=G5)
    T(s7, d, x+0.1, BY+1.34, 1.6, 0.2, sz=7.5, color=dc)

# Compliance bars
T(s7, "DOCUMENT COVERAGE BY TYPE", 0.38, BY+1.8, 5, 0.2, sz=7, bold=True, color=G5)
cov_bars = [("Salary Slips",99.2,GR),("Form 16",96,CY),("Offer Letters",91,IN),("PF Statements",87,AM),("Relieving Letters",78,RE)]
for i,(label,pct,col) in enumerate(cov_bars):
    bar(s7, 0.35, BY+2.08+i*0.46, 4.2, label, pct, f"{pct}%", col)

# Dept health
T(s7, "VAULT HEALTH BY DEPARTMENT", 0.38, BY+4.42, 5, 0.2, sz=7, bold=True, color=G5)
dept_bars = [("Finance",93,PU),("HR",91,PU),("Engineering",86,PU),("Sales",74,AM),("Customer Success",69,RE)]
for i,(label,pct,col) in enumerate(dept_bars):
    bar(s7, 0.35, BY+4.7+i*0.32, 4.0, label, pct, f"{pct}%", col)

# Alerts
R(s7, 0.22, BY+6.04, 7.7, 0.21, fill=A2)
T(s7, "  Warning: 4 exceptions unresolved — 2 past 4-hr SLA. Oldest: 31 hrs. Action needed.", 0.32, BY+6.06, 7.4, 0.18, sz=8, color=RGBColor(0x92,0x40,0x0E))

# Right: capabilities
caps7 = [
    ("Compliance Calendar","Statutory obligations tracked automatically against deadlines. Form 16 due date, PF annual statement, gratuity nominations — RAG status: Green / Amber / Red."),
    ("1-Click Audit Export","Signed, timestamped, tamper-evident PDF. Any auditor. Any period. Any document type. Backed by immutable audit_event log.","0 manual steps"),
    ("Weekly Digest — Push, Not Pull","Automated CHRO digest with vault health, doc coverage, exception queue, and alumni activity. No analyst required. No dashboard login needed.","0 manual steps"),
    ("Natural Language Analytics","Ask: 'What is our Form 16 compliance vs last quarter?' — PRANA answers from verified document data, not HR system self-reports.","Ground truth"),
    ("Statutory SLA Monitoring","PF return, ESI filing, TDS deposit — deadlines tracked per statutory calendar. Breach alerts before the due date, not after the penalty.","Auto-tracked"),
]
for i,(t,b,*extra) in enumerate(caps7):
    y = BY + i*1.23
    R(s7, 8.15, y, 4.95, 1.12, fill=IL)
    R(s7, 8.15, y, 0.03, 1.12, fill=PU)
    T(s7, t, 8.28, y+0.08, 4.6, 0.24, sz=9.5, bold=True, color=IN)
    T(s7, b, 8.28, y+0.36, 4.6, 0.65, sz=8.5, color=G5)
    if extra:
        R(s7, 11.3, y+0.1, 1.68, 0.22, fill=IL)
        T(s7, extra[0], 11.3, y+0.11, 1.68, 0.2, sz=7.5, bold=True, color=PU, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════
# SLIDE 8 — CFO
# ═══════════════════════════════════════════
s8 = prs.slides.add_slide(BLANK)
bg(s8)
header(s8, "Payroll intelligence from actual documents — not HR's input, not the payroll report.", "FOR THE CFO")
footer(s8, 8)

R(s8, 0.22, BY, 7.7, 6.25, fill=WH)
R(s8, 0.22, BY, 7.7, 0.025, fill=CY)
T(s8, "CFO Weekly Cost & Compliance Pulse  ·  Week 32", 0.38, BY+0.1, 7.3, 0.28, sz=10, bold=True, color=SL)
T(s8, "TechServe India  ·  2,000 employees  ·  Consent coverage: 97.4%", 0.38, BY+0.4, 7.3, 0.2, sz=8, color=G5)

cfo_stats = [("1,997","Headcount","Budget 2,050  −53",RE),("99.4%","TDS alignment","Salary slip vs Form 16",GR),("97.8%","CTC variance","within ±12% tolerance",GR),("₹9L","Exit cost","6 exits · repl. cost",AM)]
for i,(v,l,d,dc) in enumerate(cfo_stats):
    x = 0.35+i*1.88
    R(s8, x, BY+0.66, 1.76, 1.0, fill=BG)
    T(s8, v, x+0.1, BY+0.72, 1.6, 0.42, sz=20, bold=True, color=SL)
    T(s8, l, x+0.1, BY+1.14, 1.6, 0.2, sz=7.5, color=G5)
    T(s8, d, x+0.1, BY+1.34, 1.6, 0.2, sz=7.5, color=dc)

T(s8, "FINANCIAL DOC COMPLIANCE", 0.38, BY+1.8, 5, 0.2, sz=7, bold=True, color=G5)
fin_bars = [("Salary Slips",98,CY),("Form 16",96,CY),("PF Annual Statement",81,AM),("Appointment Letters",94,CY)]
for i,(label,pct,col) in enumerate(fin_bars):
    bar(s8, 0.35, BY+2.08+i*0.46, 4.2, label, pct, f"{pct}%", col)

T(s8, "PAYROLL INTELLIGENCE — VERIFIED FROM DOCUMENTS", 0.38, BY+4.12, 6, 0.2, sz=7, bold=True, color=G5)
R(s8, 0.35, BY+4.38, 7.35, 0.02, fill=BG)
intel_rows = [
    ("TDS consistency","99.4% monthly salary TDS aligns with annual Form 16 TDS — verified cross-document","✓",GR),
    ("CTC variance","97.8% employees within ±12% appointment letter vs actual gross salary tolerance","✓",GR),
    ("Budget utilization","₹4.4 Cr  (95.9% used, ₹18L under budget) — from verified payroll documents","⚡",AM),
    ("Attrition cost","₹33L estimated monthly · ₹9L/exit · verified tenure from actual offer + relieving letters","⚠",RE),
]
for i,(label,detail,icon,col) in enumerate(intel_rows):
    y = BY+4.44+i*0.42
    T(s8, f"{icon}  {label}", 0.38, y, 2.2, 0.22, sz=8.5, bold=True, color=col)
    T(s8, detail, 2.68, y, 5.1, 0.32, sz=8.5, color=G5)

R(s8, 0.22, BY+6.04, 7.7, 0.21, fill=A2)
T(s8, "  2 payroll anomalies — 47 employees, ₹2.3L delta. Awaiting CFO ack. SLA: 48 hrs.", 0.32, BY+6.06, 7.4, 0.18, sz=8, color=RGBColor(0x92,0x40,0x0E))

caps8 = [
    ("Consent-First Aggregation","Individual salary never visible. Minimum 30-person cohort required before any aggregate shown. DPDP-compliant by design.","Min 30"),
    ("₹50 Cr Penalty — Structurally Impossible","DPDP max penalty for data misuse is ₹250 Cr per incident. Consent-first design + per-DEK encryption makes financial data misuse architecturally impossible.","₹250Cr cap"),
    ("Industry Benchmarks from Verified Slips","Compensation benchmarks sourced from actual verified salary slips — not consultant surveys. Role / grade / location cohorts. Always consent-covered.","Verified"),
    ("Attrition Cost to the Rupee","Verified tenure from appointment + relieving letters. Replacement cost by department and grade. No spreadsheet estimates.","Ground truth"),
    ("Anomaly Ack Workflow","PAN duplicates, bulk salary revisions, CTC inconsistencies — flagged automatically. CFO one-click acknowledgement. Temporal-backed audit trail.","1-click ack"),
]
for i,(t,b,*extra) in enumerate(caps8):
    y = BY + i*1.23
    R(s8, 8.15, y, 4.95, 1.12, fill=C2)
    R(s8, 8.15, y, 0.03, 1.12, fill=CY)
    T(s8, t, 8.28, y+0.08, 4.6, 0.24, sz=9.5, bold=True, color=CY)
    T(s8, b, 8.28, y+0.36, 4.6, 0.65, sz=8.5, color=G5)
    if extra:
        R(s8, 11.35, y+0.08, 1.63, 0.22, fill=C2)
        T(s8, extra[0], 11.35, y+0.09, 1.63, 0.2, sz=7.5, bold=True, color=CY, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════
# SLIDE 9 — CISO
# ═══════════════════════════════════════════
s9 = prs.slides.add_slide(BLANK)
bg(s9)
header(s9, "Every access logged. Every document watermarked. Zero plaintext PAN — ever.", "FOR THE CISO")
footer(s9, 9)

R(s9, 0.22, BY, 7.7, 6.25, fill=WH)
R(s9, 0.22, BY, 7.7, 0.025, fill=GR)
T(s9, "InfoSec Weekly Security Pulse  ·  Week 32", 0.38, BY+0.1, 7.3, 0.28, sz=10, bold=True, color=SL)
T(s9, "TechServe India  ·  2,000 employees  ·  PRANA live 7 months", 0.38, BY+0.4, 7.3, 0.2, sz=8, color=G5)

ciso_stats = [("1,847","Doc accesses","100% watermarked",GR),("3","Anomalies","2 open · 1 auto-resolved",RE),("1","Force-logout","Bulk · Rajkot IP",RE),("34","Share tokens","31 expired · 3 active",GR)]
for i,(v,l,d,dc) in enumerate(ciso_stats):
    x = 0.35+i*1.88
    R(s9, x, BY+0.66, 1.76, 1.0, fill=BG)
    T(s9, v, x+0.1, BY+0.72, 1.6, 0.42, sz=20, bold=True, color=SL)
    T(s9, l, x+0.1, BY+1.14, 1.6, 0.2, sz=7.5, color=G5)
    T(s9, d, x+0.1, BY+1.34, 1.6, 0.2, sz=7.5, color=dc)

T(s9, "ANOMALY BREAKDOWN — THIS WEEK", 0.38, BY+1.8, 5, 0.2, sz=7, bold=True, color=G5)
anom_bars = [("Bulk access (HIGH)",36,RE),("Foreign IP (MED)",27,AM),("Failed TOTP ×4 (LOW)",27,GR),("Unusual hours",10,G5)]
for i,(label,pct,col) in enumerate(anom_bars):
    bar(s9, 0.35, BY+2.08+i*0.42, 3.8, label, pct, f"{pct}%", col)

T(s9, "INCIDENTS — THIS WEEK", 0.38, BY+3.9, 4, 0.2, sz=7, bold=True, color=G5)
incidents = [
    ("HIGH","Bulk access","47 docs in 4 min, Rajkot IP 103.21.x.x. Force-logout in 4.2 min. Watermark trail preserved. No exfil confirmed.","Resolved"),
    ("MED","Foreign IP login","OA-Operator logged in from Singapore IP 103.8.x.x. Monitoring active. CISO sees full IP; employee sees city only.","Open"),
    ("LOW","Failed TOTP ×4","PA account — auto-locked 30 min. No breach confirmed. Audit logged.","Resolved"),
]
sev_col = {"HIGH":RE,"MED":AM,"LOW":GR}
for i,(sev,label,detail,status) in enumerate(incidents):
    y = BY+4.18+i*0.64
    R(s9, 0.35, y, 0.48, 0.22, fill=sev_col[sev])
    T(s9, sev, 0.35, y+0.02, 0.48, 0.2, sz=7, bold=True, color=WH, align=PP_ALIGN.CENTER)
    T(s9, label, 0.9, y, 1.8, 0.22, sz=8.5, bold=True, color=SL)
    T(s9, detail, 0.9, y+0.24, 5.8, 0.32, sz=8, color=G5)
    sc = GR if status=="Resolved" else AM
    R(s9, 6.15, y, 0.7, 0.22, fill=LG if status=="Resolved" else A2)
    T(s9, status, 6.15, y+0.02, 0.7, 0.2, sz=7.5, bold=True, color=sc, align=PP_ALIGN.CENTER)

R(s9, 0.22, BY+6.04, 7.7, 0.21, fill=LG)
T(s9, "  Q2 posture score: 94%  ·  Avg response time: 4.2 min  ·  Zero PAN in any log  ·  100% DPDP compliant", 0.32, BY+6.06, 7.4, 0.18, sz=8, color=GR)

caps9 = [
    ("0 Plaintext PAN — Architecture, Not Policy","PAN → HMAC-SHA256 in 2ms. Plaintext destroyed. No PAN in DB row, cache entry, API response, or log line. Zero exceptions.","2ms"),
    ("Per-Employee DEK","Each employee has own encryption key. Full database compromise + no KMS access = zero readable document data.","Blast radius: 1"),
    ("RLS at Database Level","Portal Admin has zero SELECT on document rows. PostgreSQL RLS refuses the query — not application-layer logic. Database law.","DB enforced"),
    ("Immutable Audit Log","No UPDATE, no DELETE on audit_event — ever. Append-only. Tamper-evident. 7-year retention. Legal hold support.","Append-only"),
    ("LLM Isolation","Extraction LLM and analytics LLM: separate instances, separate keys, separate audit logs. No cross-contamination path.","Isolated"),
]
for i,(t,b,*extra) in enumerate(caps9):
    y = BY + i*1.23
    R(s9, 8.15, y, 4.95, 1.12, fill=LG)
    R(s9, 8.15, y, 0.03, 1.12, fill=GR)
    T(s9, t, 8.28, y+0.08, 4.6, 0.24, sz=9.5, bold=True, color=GR)
    T(s9, b, 8.28, y+0.36, 4.6, 0.65, sz=8.5, color=G5)
    if extra:
        R(s9, 11.35, y+0.08, 1.63, 0.22, fill=LG)
        T(s9, extra[0], 11.35, y+0.09, 1.63, 0.2, sz=7.5, bold=True, color=GR, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════
# SLIDE 10 — EMPLOYEE
# ═══════════════════════════════════════════
s10 = prs.slides.add_slide(BLANK)
bg(s10)
header(s10, "Your vault. Your career. Your control. Forever.", "FOR THE EMPLOYEE")
footer(s10, 10)

# Left col
R(s10, 0.22, BY, 7.7, 6.25, fill=WH)
R(s10, 0.22, BY, 7.7, 0.025, fill=VI)

T(s10, "Three promises that change everything", 0.38, BY+0.1, 7.3, 0.26, sz=9.5, bold=True, color=VI)
promises = [
    ("Cryptographically Yours","Once a document reaches your vault, your employer cannot touch it. Not open it. Not delete it. Not share it. This is not a policy. It is how the mathematics works."),
    ("Verified, Not Just Stored","Every document was pushed by your employer, extracted by AI, and timestamped immutably. When you share a salary slip, the recipient knows it is real — the system proves it. QR scan. Verified. Done."),
    ("Follows You, Forever","One vault. Every employer. Your entire career automatically assembled — never lost, always accessible — whether you are employed, between jobs, or 10 employers in."),
]
for i,(t,b) in enumerate(promises):
    y = BY+0.5+i*1.1
    R(s10, 0.35, y, 0.26, 0.26, fill=VI)
    T(s10, str(i+1), 0.35, y+0.02, 0.26, 0.24, sz=9, bold=True, color=WH, align=PP_ALIGN.CENTER)
    T(s10, t, 0.72, y+0.02, 6.9, 0.26, sz=9.5, bold=True, color=VI)
    T(s10, b, 0.72, y+0.32, 6.9, 0.72, sz=8.5, color=G5)

# Key metrics row
T(s10, "WHAT THIS MEANS IN NUMBERS", 0.38, BY+3.9, 6, 0.2, sz=7, bold=True, color=G5)
emp_metrics = [("2ms","PAN destroyed forever"),("₹0","Cost to verify for a bank"),("0","Employer accesses post-exit"),("38s","Alumni vault access time")]
for i,(v,l) in enumerate(emp_metrics):
    x = 0.35+i*1.88
    R(s10, x, BY+4.18, 1.76, 0.82, fill=BG)
    T(s10, v, x+0.1, BY+4.24, 1.58, 0.36, sz=18, bold=True, color=VI)
    T(s10, l, x+0.1, BY+4.6, 1.58, 0.32, sz=7.5, color=G5)

# DPDP rights
T(s10, "YOUR RIGHTS UNDER DPDP ACT 2023 — ALL BUILT IN", 0.38, BY+5.12, 6, 0.2, sz=7, bold=True, color=G5)
R(s10, 0.35, BY+5.38, 7.35, 0.02, fill=BG)
rights10 = [("S.11","Right to Access","See everything PRANA holds. Download it all — one click."),("S.12","Right to Correction + Erasure","Flag errors. Request deletion with DPDP certificate."),("S.13","Right to Grievance Redressal","Track grievances. One-click DPB escalation."),("S.14","Right to Nomination","Nominate family for vault access. First platform in India."),("S.7","Right to Withdraw Consent","Analytics toggles. Immediate effect. No fine print.")]
for i,(sec,title,detail) in enumerate(rights10):
    y = BY+5.44+i*0.35
    R(s10, 0.35, y, 0.4, 0.22, fill=IL)
    T(s10, sec, 0.35, y+0.02, 0.4, 0.2, sz=7, bold=True, color=PU, align=PP_ALIGN.CENTER)
    T(s10, title, 0.85, y, 2.5, 0.22, sz=8.5, bold=True, color=SL)
    T(s10, detail, 3.45, y, 4.1, 0.22, sz=8, color=G5)

# Right: capabilities
caps10 = [
    ("My Vault — All Docs, All Employers","Salary slips, Form 16, offer letters, appointment letters, relieving letters, PF statements — one place, all employers, your career lifetime."),
    ("C-Share — Verified Document Sharing","Named recipient. Time-limited token. OTP-gated. Watermarked. View-count visible. Revocable any time. Bank scans QR — verified without HR."),
    ("Activity Log — No Black Boxes","Every actor who touched your document: name, timestamp, IP, channel. Immutable. You see it. Employer cannot hide it."),
    ("Document Alerts","Missing salary slip for a month? PRANA tells you. Relieving letter overdue? One-click tracked request to employer. Proactive, not reactive."),
    ("Mobile App — Face ID + Biometric","Face ID / fingerprint re-auth for every document access. Server-side watermark before bytes reach device. Secure storage: SecureStore, not AsyncStorage."),
]
for i,(t,b) in enumerate(caps10):
    y = BY + i*1.23
    R(s10, 8.15, y, 4.95, 1.12, fill=V2)
    R(s10, 8.15, y, 0.03, 1.12, fill=VI)
    T(s10, t, 8.28, y+0.08, 4.6, 0.26, sz=9.5, bold=True, color=VI)
    T(s10, b, 8.28, y+0.38, 4.6, 0.65, sz=8.5, color=G5)

# ═══════════════════════════════════════════
# SLIDE 11 — AI PIPELINE
# ═══════════════════════════════════════════
s11 = prs.slides.add_slide(BLANK)
bg(s11)
header(s11, "Local LLM. Full document in. Insights only out. Data never leaves India.", "AI PIPELINE + GUARDRAILS")
footer(s11, 11)

# 6 stages
stages = [
    ("01","INGEST","Validate · S3 put\n1 DB write · Kafka\n202 in <50ms",IL),
    ("02","OCR","Tesseract (local)\n→ Textract (fallback)\nMultilingual",LG),
    ("03","LLM EXTRACT","Qwen 2.5-14B\nFull doc in\nStructured JSON out",IL),
    ("04","RESOLVE","PAN dedup · EPFO\nIdentity match\nCross-tenant",LG),
    ("05","INSIGHTS","Llama 3.1-8B\nInsights only out\nPrivacy filter ENFORCED",IL),
    ("06","ROUTE","Vault routing\nSSE push\nEmployee access",LG),
]
bw = 2.08
for i,(num,title,body,bg_c) in enumerate(stages):
    x = 0.22+i*(bw+0.04)
    R(s11, x, BY, bw, 2.4, fill=bg_c)
    R(s11, x, BY, bw, 0.025, fill=PU if i%2==0 else GR)
    R(s11, x, BY, 0.025, 2.4, fill=PU if i%2==0 else GR)
    R(s11, x+0.1, BY+0.1, 0.3, 0.3, fill=PU if i%2==0 else GR)
    T(s11, num, x+0.1, BY+0.1, 0.3, 0.3, sz=8.5, bold=True, color=WH, align=PP_ALIGN.CENTER)
    T(s11, title, x+0.5, BY+0.12, bw-0.6, 0.28, sz=9, bold=True, color=SL)
    T(s11, body, x+0.12, BY+0.52, bw-0.2, 1.75, sz=8.5, color=G5)
    if i < 5:
        T(s11, ">", x+bw+0.01, BY+1.0, 0.08, 0.4, sz=14, color=G5)

# LLM models detail box
R(s11, 0.22, BY+2.55, 8.5, 1.65, fill=DK)
R(s11, 0.22, BY+2.55, 8.5, 0.025, fill=PU)
T(s11, "LLM STACK — ALL LOCALLY HOSTED, INSIDE YOUR INFRASTRUCTURE", 0.38, BY+2.65, 7.8, 0.22, sz=8, bold=True, color=LA)
llm_rows = [
    ("Extraction","Qwen/Qwen2.5-14B-Instruct","Best structured JSON for Indian HR documents. Salary slips, Form 16, offer letters. Hindi + English."),
    ("Insights / RAG","meta-llama/Llama-3.1-8B-Instruct","Career insights, chatbot answers. Privacy filter runs AFTER this model. Raw figures blocked."),
    ("Embeddings","BAAI/bge-m3","Multilingual — Hindi, English, Marathi. 512-token chunks, 50-token overlap. Qdrant vector store."),
    ("OCR","Tesseract (local) + AWS Textract","Tesseract primary. Textract fallback. No document bytes leave PRANA infrastructure to Textract without org consent."),
]
for i,(role,model,detail) in enumerate(llm_rows):
    y = BY+2.95+i*0.3
    T(s11, role, 0.38, y, 1.2, 0.22, sz=8, bold=True, color=LA)
    T(s11, model, 1.68, y, 3.2, 0.22, sz=8, bold=True, color=WH)
    T(s11, detail, 5.0, y, 3.55, 0.22, sz=7.5, color=RGBColor(0x94,0xA3,0xB8))

# Privacy filter + data sovereignty
R(s11, 8.95, BY+2.55, 4.15, 1.65, fill=R2)
R(s11, 8.95, BY+2.55, 4.15, 0.025, fill=RE)
T(s11, "PRIVACY FILTER — NON-NEGOTIABLE", 9.08, BY+2.65, 3.9, 0.24, sz=9, bold=True, color=RE)
T(s11, "Stage 05 output passes mandatory filter.\nAny raw salary figure, PAN, or NIK → BLOCKED.\n\nResponse becomes:\n\"I can share insights about your career but cannot share specific financial figures.\"\n\nHardcoded. Not configurable. Not bypassable by any admin.", 9.08, BY+2.95, 3.85, 1.15, sz=8.5, color=RE)

# Data sovereignty + capabilities
R(s11, 0.22, BY+4.38, 12.88, 1.98, fill=WH)
R(s11, 0.22, BY+4.38, 12.88, 0.025, fill=IL)
sov_items = [
    ("No Document Leaves India","HuggingFace today → Ollama / vLLM on-premise tomorrow. Full data sovereignty. Zero calls to OpenAI, Anthropic, or any external LLM API."),
    ("Multilingual from Day 1","BGE-M3 embeddings handle Hindi salary slips, English offer letters, Marathi PF statements — same pipeline, same vault."),
    ("LLM Isolation","Extraction LLM and analytics LLM: separate GPU instances, separate keys, separate audit logs. No cross-contamination path architecturally possible."),
    ("90+ System Notifications","AI events (stage completion, exception flagged, anomaly detected) trigger 90+ automated notifications across Email, SMS, WhatsApp, Push, Portal Bell, Incident Log."),
]
for i,(t,b) in enumerate(sov_items):
    x = 0.35+i*3.25
    T(s11, t, x, BY+4.52, 3.0, 0.26, sz=9, bold=True, color=IN)
    T(s11, b, x, BY+4.82, 3.0, 1.38, sz=8, color=G5)

# ═══════════════════════════════════════════
# SLIDE 12 — LABOUR LAW
# ═══════════════════════════════════════════
s12 = prs.slides.add_slide(BLANK)
bg(s12)
header(s12, "29 Acts. 4 Labour Codes. One vault that covers the full compliance stack.", "LABOUR LAW COMPLIANCE")
footer(s12, 12)

codes = [
    ("CODE ON WAGES","Wage slips (monthly)\nWage register export\nBonus records\n3-year mandatory retention lock\nMinimum wage variance tracking",IN,IL),
    ("INDUSTRIAL RELATIONS","Appointment letters\nRetrenchment notices\nGrievance records + SLA tracking\nStanding orders compliance\nIncrease / promotion letter trail",VI,V2),
    ("SOCIAL SECURITY CODE","PF annual statement (EPFO)\nESI card + contribution records\nGratuity Forms F / I / L / M\nMaternity records (dual-consent)\nSuperannuation documents",GR,LG),
    ("OSH CODE","Contract labour records\nMigrant worker passbook\nHealth certificate tracking\nSafety training records\nFactory / establishment register",CY,C2),
]
for i,(title,body,ac,bg_c) in enumerate(codes):
    x = 0.22+i*3.28
    R(s12, x, BY, 3.1, 4.1, fill=bg_c)
    R(s12, x, BY, 3.1, 0.025, fill=ac)
    R(s12, x, BY, 0.025, 4.1, fill=ac)
    T(s12, title, x+0.12, BY+0.1, 2.85, 0.24, sz=8.5, bold=True, color=ac)
    T(s12, body, x+0.12, BY+0.45, 2.85, 3.5, sz=8.5, color=SL)

# Critical gaps closed
R(s12, 0.22, BY+4.28, 12.88, 2.08, fill=WH)
R(s12, 0.22, BY+4.28, 12.88, 0.025, fill=IL)
T(s12, "CRITICAL GAPS CLOSED BY PRANA", 0.38, BY+4.38, 8, 0.22, sz=8, bold=True, color=G5)

gaps = [
    ("Gratuity Form F / I / L / M","All four now in PRANA taxonomy — nomination, application, payment, rejection. Timestamped. Disputes now have an immutable paper trail. Previously zero platforms tracked all four."),
    ("Form 16 — Repository, Not Issuer","PRANA stores employer-pushed Parts A + B. TRACES is authoritative. PRANA is the employee's permanent, portable copy — accessible post-exit, shareable with bank or CA. Zero re-issuance needed."),
    ("Maternity Records — Sensitive Category","Enhanced privacy flag per PRANA taxonomy. Dual-consent required for any share or analytics inclusion. Medical-adjacent data treated with strictest access controls."),
    ("IT Act + Income Tax Act Coverage","Form 16 TDS compliance, salary TDS consistency (99.4% alignment verified), investment proofs — covered under same vault alongside Labour Code obligations."),
]
for i,(t,b) in enumerate(gaps):
    x = 0.35+i*3.25
    T(s12, t, x, BY+4.68, 3.0, 0.28, sz=9, bold=True, color=IN)
    T(s12, b, x, BY+5.0, 3.0, 1.28, sz=8, color=G5)

# ═══════════════════════════════════════════
# SLIDE 13 — DPDP ACT 2023
# ═══════════════════════════════════════════
s13 = prs.slides.add_slide(BLANK)
bg(s13)
header(s13, "Employee data rights — built in at the architecture level, not added as a feature.", "DPDP ACT 2023")
footer(s13, 13)

# 5 rights across top
rights13 = [
    ("S.11","Right to Access","See all data PRANA holds. Download it all. One click.","30d / instant",GR),
    ("S.12","Right to Correction + Erasure","Flag errors. Request deletion with DPDP certificate. Audit logs survive — 7yr legal retention explained.","30d erasure",PU),
    ("S.13","Right to Grievance Redressal","Track grievances. Auto-ack within 48 hrs. One-click DPB escalation if unresolved in 30 days.","48h ack",AM),
    ("S.14","Right to Nomination","Nominate family member for vault access on incapacitation. First platform in India to implement this right.","Instant",CY),
    ("S.7","Right to Withdraw Consent","Per-purpose consent toggles. Withdrawal immediate — Kafka notifies all services within milliseconds. No fine print.","Instant",RE),
]
rw = 2.55
for i,(sec,title,body,sla,col) in enumerate(rights13):
    x = 0.22+i*(rw+0.04)
    R(s13, x, BY, rw, 2.65, fill=WH)
    R(s13, x, BY, rw, 0.025, fill=col)
    R(s13, x, BY, 0.025, 2.65, fill=col)
    R(s13, x+0.1, BY+0.1, 0.38, 0.22, fill=col)
    T(s13, sec, x+0.1, BY+0.11, 0.38, 0.2, sz=7.5, bold=True, color=WH, align=PP_ALIGN.CENTER)
    T(s13, title, x+0.6, BY+0.1, rw-0.7, 0.28, sz=9, bold=True, color=col)
    T(s13, body, x+0.12, BY+0.48, rw-0.2, 1.85, sz=8.5, color=G5)
    R(s13, x+0.12, BY+2.36, 0.8, 0.2, fill=col)
    T(s13, sla, x+0.12, BY+2.37, 0.8, 0.2, sz=7.5, bold=True, color=WH, align=PP_ALIGN.CENTER)

# Consent model + Q2 compliance
R(s13, 0.22, BY+2.82, 6.3, 3.56, fill=DK)
R(s13, 0.22, BY+2.82, 6.3, 0.025, fill=PU)
T(s13, "CONSENT MODEL", 0.38, BY+2.92, 5.8, 0.24, sz=10, bold=True, color=LA)
consent_pts = [
    "Consent is per-tenant, per-purpose — not a global toggle",
    "Every grant and withdrawal timestamped in consent_log table",
    "Withdrawal is immediate — Kafka notifies all downstream services in <1 sec",
    "Aggregated analytics require minimum 30-person cohort consent coverage",
    "Audit logs survive erasure: 7-year legal retention overrides DPDP S.12",
    "PRANA explains this to employees proactively — not buried in policy",
    "100% India data residency — verified, zero cross-border transfers",
    "Max DPDP penalty exposure: ₹250 Cr — misuse structurally impossible by architecture",
]
for i,pt in enumerate(consent_pts):
    y = BY+3.24+i*0.38
    R(s13, 0.38, y+0.07, 0.07, 0.07, fill=LA)
    T(s13, pt, 0.55, y, 5.75, 0.28, sz=8.5, color=WH)

# Right: Q2 compliance record + Temporal
R(s13, 6.75, BY+2.82, 6.35, 1.65, fill=LG)
R(s13, 6.75, BY+2.82, 6.35, 0.025, fill=GR)
T(s13, "Q2 COMPLIANCE RECORD", 6.9, BY+2.92, 5.9, 0.24, sz=10, bold=True, color=GR)
q2 = [("100%","All requests within SLA","Q2 total: 5 DPDP requests"),("1","Erasure completed","Day 12 of 30 — well within SLA"),("1","Export completed","Day 3 of 3 — well within SLA"),("100%","Audit trail complete","All workflows logged to audit_event")]
for i,(v,l,d) in enumerate(q2):
    x = 6.9+i*1.55
    T(s13, v, x, BY+3.28, 1.45, 0.38, sz=18, bold=True, color=GR)
    T(s13, l, x, BY+3.65, 1.45, 0.2, sz=7.5, bold=True, color=SL)
    T(s13, d, x, BY+3.85, 1.45, 0.2, sz=7.5, color=G5)

R(s13, 6.75, BY+4.65, 6.35, 1.72, fill=IL)
R(s13, 6.75, BY+4.65, 6.35, 0.025, fill=PU)
T(s13, "TEMPORAL WORKFLOWS — NOT MANUAL PROCESSES", 6.9, BY+4.75, 5.9, 0.24, sz=9, bold=True, color=IN)
wf_list = [("ErasureWorkflow","Durable, retryable — survives pod restarts. SLA timer built-in. Escalates to PA on day 28."),("DataExportWorkflow","Packages all data, generates DPDP certificate, emails to employee. Audit-logged."),("ConsentWithdrawalWorkflow","Immediate Kafka fan-out. All services notified. Confirmed by consumer receipts."),("GrievanceWorkflow","48-hr auto-ack. Day-30 DPB escalation if unresolved. All timestamped.")]
for i,(wf,detail) in enumerate(wf_list):
    y = BY+5.06+i*0.32
    T(s13, wf, 6.9, y, 2.3, 0.22, sz=8.5, bold=True, color=PU)
    T(s13, detail, 9.32, y, 3.6, 0.28, sz=8, color=G5)

# ═══════════════════════════════════════════
# SLIDE 14 — SECURITY + SCALE
# ═══════════════════════════════════════════
s14 = prs.slides.add_slide(BLANK)
bg(s14)
header(s14, "Enterprise-grade encryption. India-scale infrastructure. 99.99% uptime.", "SECURITY + SCALE")
footer(s14, 14)

# Encryption model
R(s14, 0.22, BY, 6.3, 2.65, fill=WH)
R(s14, 0.22, BY, 6.3, 0.025, fill=PU)
T(s14, "THREE-KEY ENVELOPE ENCRYPTION", 0.38, BY+0.1, 5.8, 0.24, sz=9, bold=True, color=IN)
enc_flow = [("Document",""), ("→ Employee DEK","FF3-1 FPE for PAN"),("→ Tenant KEK","Per-org key"),("→ AWS KMS","ap-south-1")]
x_e = 0.38
for i,(label,sub) in enumerate(enc_flow):
    bg_e = [RGBColor(0xFE,0xF3,0xC7),IL,V2,LG][i] if i>0 else BG
    R(s14, x_e, BY+0.48, 1.35, 0.62, fill=bg_e)
    T(s14, label, x_e+0.06, BY+0.54, 1.24, 0.24, sz=8, bold=True, color=SL)
    T(s14, sub, x_e+0.06, BY+0.76, 1.24, 0.2, sz=7, color=G5)
    x_e += 1.42

enc_details = [
    ("Passwords","Argon2id  ·  time=2  ·  memory=65,536  ·  parallelism=2"),
    ("PAN","FF3-1 Format-Preserving Encryption per employee DEK  ·  HMAC dedup key"),
    ("TOTP secret","AES-256-GCM  ·  Never stored in plaintext, never in logs"),
    ("JWT revocation","Redis namespace revoked:{jti}  ·  TTL = session lifetime"),
    ("Cache keys","pan_token (HMAC-SHA256) ONLY  ·  Zero plaintext PAN in any cache entry"),
]
for i,(label,detail) in enumerate(enc_details):
    y = BY+1.2+i*0.28
    T(s14, label, 0.38, y, 1.3, 0.22, sz=8.5, bold=True, color=PU)
    T(s14, detail, 1.78, y, 4.6, 0.22, sz=8.5, color=G5)

# Scale stats
R(s14, 6.75, BY, 6.35, 2.65, fill=DK)
R(s14, 6.75, BY, 6.35, 0.025, fill=LA)
T(s14, "BUILT FOR INDIA'S ENTIRE WORKFORCE", 6.9, BY+0.1, 5.9, 0.24, sz=9, bold=True, color=LA)
scale_m = [("1 Lakh+","Organisations"),("1 Crore+","Employees"),("125 PB","Document storage"),("53","Temporal workflows"),("5","Kafka topics (12 partitions each)"),("7","Year audit retention")]
for i,(v,l) in enumerate(scale_m):
    row,col = divmod(i,3)
    x = 6.9+col*2.1; y = BY+0.5+row*1.0
    T(s14, v, x, y, 1.95, 0.45, sz=22, bold=True, color=LA)
    T(s14, l, x, y+0.45, 1.95, 0.24, sz=8, color=RGBColor(0x94,0xA3,0xB8))

# Infra row
R(s14, 0.22, BY+2.82, 12.88, 3.56, fill=WH)
R(s14, 0.22, BY+2.82, 12.88, 0.025, fill=IL)
T(s14, "INFRASTRUCTURE STACK", 0.38, BY+2.92, 5, 0.22, sz=8, bold=True, color=G5)
infra = [
    ("YugabyteDB","Dual-region: Mumbai + Hyderabad  ·  256 tablets  ·  CRDT active-active\n99.99% uptime SLA  ·  PostgreSQL-compatible  ·  RLS enforced","Database"),
    ("Apache Kafka","AWS MSK  ·  KRaft mode  ·  5 topics, 12 partitions each\nMirrorMaker 2 bidirectional cross-region sync  ·  5 typed consumers","Event Stream"),
    ("Redis CRDT","ElastiCache Global Datastore  ·  Active-active both regions\nSub-10ms cross-region sync  ·  4 namespaces: identity, share, vault, JWT","Cache"),
    ("Temporal","Python SDK v1.x  ·  53 named workflows  ·  5 task queues\nZero cron, zero Celery  ·  Thin workflow shell + rich service class pattern","Workflow"),
    ("AWS KMS","ap-south-1 (primary)  ·  Customer-managed keys\nTenant KEK wraps employee DEK  ·  Platform secret for HMAC pan_token","Key Mgmt"),
    ("AWS S3 / MinIO","Document storage  ·  Server-side encryption  ·  Lifecycle policies\nDual-region replication  ·  125 PB capacity design","Storage"),
]
for i,(name,detail,badge) in enumerate(infra):
    row,col = divmod(i,3)
    x = 0.35+col*4.32; y = BY+3.22+row*1.58
    R(s14, x, y, 4.12, 1.42, fill=G2)
    R(s14, x, y, 0.025, 1.42, fill=PU)
    R(s14, x+0.1, y+0.08, 0.5, 0.22, fill=BG)
    T(s14, badge, x+0.1, y+0.09, 0.5, 0.2, sz=6.5, bold=True, color=G5, align=PP_ALIGN.CENTER)
    T(s14, name, x+0.72, y+0.08, 3.3, 0.26, sz=10, bold=True, color=SL)
    T(s14, detail, x+0.12, y+0.42, 3.85, 0.9, sz=8, color=G5)

# ═══════════════════════════════════════════
# SLIDE 15 — IMMUTABLE AUDIT + TRUST
# ═══════════════════════════════════════════
s15 = prs.slides.add_slide(BLANK)
bg(s15)
header(s15, "We cannot see your salary. We cannot access your docs. By architecture, not policy.", "IMMUTABLE AUDIT + TRUST")
footer(s15, 15)

# Left: audit chain
R(s15, 0.22, BY, 6.5, 6.25, fill=WH)
R(s15, 0.22, BY, 6.5, 0.025, fill=PU)
T(s15, "EVERY EVENT — PERMANENTLY RECORDED — APPEND-ONLY", 0.38, BY+0.1, 6.1, 0.22, sz=7.5, bold=True, color=G5)

chain15 = [
    ("PUSH","Document pushed by employer","tenant_id · document_id · actor · IP · timestamp","sha256:a3f8c2e1...","Jun 24, 09:14:22 IST",IL,PU),
    ("AI","AI pipeline: Stage 03 LLM extraction","Qwen 2.5-14B · structured JSON · insights only · no raw figures","sha256:b7d1e4f8...","Jun 24, 09:14:31 IST",LG,GR),
    ("VIEW","Employee viewed document","Mobile · Face ID re-auth · watermark applied · IP logged · channel: MOBILE","sha256:c9f2a1d3...","Jun 24, 11:02:47 IST",A2,AM),
    ("SHARE","Share token issued","Recipient: HDFC Bank Officer · OTP-gated · 10-min session · watermark: recipient+timestamp+token","sha256:d4e8b3c7...","Jun 24, 14:37:08 IST",V2,VI),
    ("REVOKE","Share token expired","Auto-expiry · access revoked · watermark trail preserved · no document bytes cached","sha256:e2f9c5a8...","Jun 24, 14:47:08 IST",C2,CY),
]
for i,(action,title,detail,h,ts,bg_c,ac) in enumerate(chain15):
    y = BY+0.4+i*1.12
    R(s15, 0.35, y, 6.22, 1.0, fill=bg_c)
    R(s15, 0.35, y, 0.5, 1.0, fill=ac)
    T(s15, action, 0.35, y+0.37, 0.5, 0.24, sz=7.5, bold=True, color=WH, align=PP_ALIGN.CENTER)
    T(s15, title, 0.95, y+0.06, 5.4, 0.26, sz=9, bold=True, color=SL)
    T(s15, detail, 0.95, y+0.36, 5.4, 0.28, sz=8, color=G5)
    T(s15, h, 0.95, y+0.68, 2.8, 0.2, sz=7.5, color=PU, italic=True)
    T(s15, ts, 4.5, y+0.06, 2.1, 0.2, sz=7.5, color=G5, align=PP_ALIGN.RIGHT)
    if i < 4:
        R(s15, 0.35, y+1.0, 0.5, 0.14, fill=BG)

# Right: trust claims
R(s15, 6.95, BY, 6.15, 2.28, fill=DK)
R(s15, 6.95, BY, 6.15, 0.025, fill=LA)
T(s15, "7-YEAR RETENTION — NON-NEGOTIABLE", 7.1, BY+0.1, 5.8, 0.24, sz=9.5, bold=True, color=LA)
T(s15, "Hot storage (0–2 years): YugabyteDB — fully queryable\nCold storage (2–7 years): Apache Iceberg on S3 — tamper-evident\nArchival: AuditArchivalWorkflow migrates automatically\n\nAudit logs survive erasure requests.\nDPDP S.12 erasure right does not override 7-year legal retention.\nPRANA explains this proactively — employee receives written explanation, not silence.", 7.1, BY+0.42, 5.8, 1.7, sz=8.5, color=WH)

# Zero access claim
R(s15, 6.95, BY+2.44, 6.15, 1.52, fill=R2)
R(s15, 6.95, BY+2.44, 6.15, 0.025, fill=RE)
T(s15, "ZERO ACCESS — BY ARCHITECTURE", 7.1, BY+2.54, 5.8, 0.26, sz=9.5, bold=True, color=RE)
za_pts = [
    "LLM output = insights only. Raw salary figures never written to any DB table.",
    "PAN never in any DB row, cache entry, log line, or API response — ever.",
    "Portal Admin has zero SELECT on document rows (PostgreSQL RLS — database law).",
    "Watermark on every share: recipient name + timestamp + token hash. Server-side. Cannot be removed.",
]
for i,pt in enumerate(za_pts):
    R(s15, 7.1, BY+2.9+i*0.26, 0.06, 0.06, fill=RE)
    T(s15, pt, 7.26, BY+2.88+i*0.26, 5.6, 0.22, sz=8.5, color=RE)

# Watermark detail
R(s15, 6.95, BY+4.12, 6.15, 2.14, fill=IL)
R(s15, 6.95, BY+4.12, 6.15, 0.025, fill=PU)
T(s15, "WATERMARK ON EVERY ACCESS", 7.1, BY+4.22, 5.8, 0.26, sz=9.5, bold=True, color=IN)
T(s15, "Content: Recipient name + access timestamp + share token (last 12 chars)\nApplied: Server-side, before bytes reach device. Client cannot bypass.\nLogged: document_access_log.watermark_applied = TRUE — mandatory field.\nPurpose: In the event of exfiltration, evidence trail always preserved.\nEmployee own access: watermarked with employee name + timestamp.", 7.1, BY+4.55, 5.8, 1.55, sz=8.5, color=G5)

# ═══════════════════════════════════════════
# SLIDE 16 — MARKET LANDSCAPE
# ═══════════════════════════════════════════
s16 = prs.slides.add_slide(BLANK)
bg(s16)
header(s16, "PRANA occupies a category no existing product covers.", "MARKET LANDSCAPE")
footer(s16, 16)

# Comparison table
R(s16, 0.22, BY, 12.88, 0.35, fill=DK)
headers16 = ["Capability","DigiLocker","HRMS Platforms\n(SAP / Darwinbox / Keka)","Document Mgmt\n(Leegality / DMS tools)","PRANA"]
col_w = [2.9, 2.2, 2.6, 2.2, 2.6]
x_t = 0.25
for i,h in enumerate(headers16):
    bg_hdr = PU if i==4 else DK
    T(s16, h, x_t, BY+0.05, col_w[i], 0.28, sz=8, bold=True, color=WH if i!=4 else LA, align=PP_ALIGN.CENTER)
    x_t += col_w[i]+0.05

rows16 = [
    ("Employer pushes documents",       "✗",     "✓ (source)",   "✗",   "✓"),
    ("Employee owns post-exit",         "✗",     "✗",             "✗",   "✓"),
    ("AI extraction + insights",        "✗",     "✗",             "✗",   "✓"),
    ("Per-employee encryption (DEK)",   "✗",     "✗",             "✗",   "✓"),
    ("Immutable audit log",             "Partial","✗",             "✗",   "✓"),
    ("DPDP Act 2023 — all 6 rights",    "✗",     "✗",             "Partial","✓"),
    ("Alumni / post-exit access",       "✗",     "✗",             "✗",   "✓"),
    ("Labour law compliance tracking",  "✗",     "Partial",        "✗",   "✓"),
    ("0 plaintext PAN in DB",           "✗",     "✗",             "✗",   "✓"),
    ("Watermarked secure share",        "✗",     "✗",             "Partial","✓"),
    ("Multilingual Hindi + English",    "✗",     "✗",             "✗",   "✓"),
    ("Portable across employers",       "Partial","✗",             "✗",   "✓"),
]
for i,row in enumerate(rows16):
    y = BY+0.4+i*0.48
    bg_r = WH if i%2==0 else G2
    R(s16, 0.22, y, 12.88, 0.46, fill=bg_r)
    x_r = 0.25
    col_colors = [SL,
                  GR if row[1]=="✓" else (AM if row[1]=="Partial" else RE),
                  GR if row[2]=="✓" else (AM if "Partial" in row[2] or row[2]=="Partial" else RE),
                  GR if row[3]=="✓" else (AM if row[3]=="Partial" else RE),
                  GR]
    for j,cell in enumerate(row):
        bold = j==0 or j==4
        T(s16, cell, x_r, y+0.12, col_w[j], 0.22, sz=8.5, bold=bold,
          color=col_colors[j] if j>0 else SL, align=PP_ALIGN.CENTER if j>0 else PP_ALIGN.LEFT)
        x_r += col_w[j]+0.05

# PRANA column highlight
R(s16, 11.42, BY, 1.68, 6.22, fill=RGBColor(0xEE,0xF2,0xFF))
R(s16, 11.42, BY, 1.68, 0.025, fill=PU)

# ═══════════════════════════════════════════
# SLIDE 17 — WHY NOW + NEXT STEP
# ═══════════════════════════════════════════
s17 = prs.slides.add_slide(BLANK)
bg(s17)
header(s17, "Three forces are converging. The window to move first is now.", "WHY NOW + NEXT STEP")
footer(s17, 17)

forces17 = [
    ("DPDP Act 2023","Regulatory Deadline","Every employer handling employee data must comply. Fines up to ₹250 Cr per incident. ₹500 Cr for repeated violation. PRANA is the fastest path — not a 12-month IT project, a SaaS subscription with 30-day onboarding.\n\nThe clock is running. Enforcement guidance expected within 2025.",RE,R2),
    ("15 Crore Workers","Portability Crisis","India's gig, contract, and formal sector workforce — the fastest-growing, most document-vulnerable segment. Average worker changes jobs every 2.8 years. Zero portable document trail. Banks, visa officers, and background agencies reject applications daily.\n\nThe problem compounds with every job change.",AM,A2),
    ("HRMS Fragmentation","Peak Pain","The average Indian enterprise runs 3–4 HR systems — SAP for corporate, Darwinbox for operations, Keka for startup entities, legacy tools for plants. Zero interoperability. PRANA ingests from all of them — no rip-and-replace, no IT committee.",PU,V2),
]
for i,(title,sub,body,ac,bg_c) in enumerate(forces17):
    x = 0.22+i*4.37
    R(s17, x, BY, 4.2, 4.6, fill=bg_c)
    R(s17, x, BY, 4.2, 0.025, fill=ac)
    R(s17, x, BY, 0.025, 4.6, fill=ac)
    R(s17, x+0.12, BY+0.1, 0.26, 0.26, fill=ac)
    T(s17, str(i+1), x+0.12, BY+0.11, 0.26, 0.24, sz=9.5, bold=True, color=WH, align=PP_ALIGN.CENTER)
    T(s17, title, x+0.48, BY+0.1, 3.6, 0.28, sz=11, bold=True, color=ac)
    R(s17, x+0.48, BY+0.42, 1.5, 0.2, fill=ac)
    T(s17, sub, x+0.48, BY+0.43, 1.5, 0.2, sz=7.5, bold=True, color=WH, align=PP_ALIGN.CENTER)
    T(s17, body, x+0.12, BY+0.76, 3.95, 3.6, sz=8.5, color=SL)

# CTA box
R(s17, 0.22, BY+4.78, 12.88, 1.58, fill=DK)
R(s17, 0.22, BY+4.78, 12.88, 0.025, fill=PU)
T(s17, "THE OFFER: Start with one department. Own your workforce document story in 30 days.", 0.42, BY+4.9, 9, 0.32, sz=12, bold=True, color=WH)
cta_pts = ["No HRMS replacement required — PRANA integrates as a push layer above existing systems","No IT committee or 6-month procurement — SaaS subscription, standard MSA, 30-day onboarding","No employee action required — documents appear in vault automatically when employer pushes","Full DPDP Act 2023 compliance on Day 1 — consent model, data rights, audit trail — all built in"]
for i,pt in enumerate(cta_pts):
    R(s17, 0.42, BY+5.38+i*0.22, 0.07, 0.07, fill=GR)
    T(s17, pt, 0.6, BY+5.34+i*0.22, 11.5, 0.22, sz=8.5, color=LA)

T(s17, "Contact: Nilesh Shah  ·  n.shah12@gmail.com  ·  prana.in", 10.2, BY+5.38, 2.7, 0.5, sz=8, color=RGBColor(0x94,0xA3,0xB8), align=PP_ALIGN.RIGHT)

# ─── SAVE ───
out = r"C:\Nilesh\claude-code\prana-docs\PRANA_CXO_Deck_v2.pptx"
prs.save(out)
print(f"Saved to: {out}")
print(f"Slides: {len(prs.slides)}")
