from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]

# ─── colours ───
NAVY   = RGBColor(0x0F, 0x0C, 0x29)
INDIGO = RGBColor(0x31, 0x2E, 0x81)
VIOLET = RGBColor(0x7C, 0x3A, 0xED)
PURPLE = RGBColor(0x63, 0x66, 0xF1)
LAVEND = RGBColor(0xA5, 0xB4, 0xFC)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
SLATE  = RGBColor(0x0F, 0x17, 0x2A)
GRAY   = RGBColor(0x64, 0x74, 0x8B)
LGRAY  = RGBColor(0xF8, 0xFA, 0xFF)
GREEN  = RGBColor(0x05, 0x96, 0x69)
LGREEN = RGBColor(0xEC, 0xFD, 0xF5)
AMBER  = RGBColor(0xF5, 0x9E, 0x0B)
RED    = RGBColor(0xE1, 0x1D, 0x48)
CYAN   = RGBColor(0x08, 0x91, 0xB2)
INDL   = RGBColor(0xEE, 0xF2, 0xFF)
BGWARM = RGBColor(0xFF, 0xFB, 0xEB)
BGRED  = RGBColor(0xFF, 0xF1, 0xF2)
BGVIO  = RGBColor(0xF5, 0xF3, 0xFF)
BGCYA  = RGBColor(0xF0, 0xF9, 0xFF)
DK     = RGBColor(0x1E, 0x1B, 0x4B)

def rect(slide, l, t, w, h, fill=None, alpha=None):
    s = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    s.line.fill.background()
    if fill:
        s.fill.solid(); s.fill.fore_color.rgb = fill
    else:
        s.fill.background()
    return s

def txt(slide, text, l, t, w, h, size=10, bold=False, color=SLATE,
        align=PP_ALIGN.LEFT, italic=False):
    txb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    txb.word_wrap = True
    tf = txb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = bold
    r.font.color.rgb = color; r.font.italic = italic
    return txb

def logo(slide):
    rect(slide, 0.18, 0.12, 0.22, 0.22, fill=INDL)
    txt(slide, "P", 0.22, 0.13, 0.15, 0.20, size=10, bold=True, color=PURPLE)
    txt(slide, "PRANA", 0.42, 0.14, 1.0, 0.20, size=10, bold=True, color=INDIGO)

def footer(slide, n, total=15):
    rect(slide, 0, 7.2, 13.33, 0.3, fill=RGBColor(0xF1,0xF5,0xF9))
    txt(slide, "Private & Confidential", 0.2, 7.22, 4, 0.25, size=7.5, color=GRAY)
    txt(slide, "NS", 6.4, 7.22, 1, 0.25, size=7.5, color=GRAY, align=PP_ALIGN.CENTER)
    txt(slide, f"Slide {n} of {total}", 9.1, 7.22, 4, 0.25, size=7.5, color=GRAY, align=PP_ALIGN.RIGHT)

def hdr(slide, badge, title, sub=""):
    rect(slide, 0, 0.5, 13.33, 1.1, fill=LGRAY)
    txt(slide, badge, 0.35, 0.54, 8, 0.25, size=8, bold=True, color=PURPLE)
    txt(slide, title, 0.35, 0.74, 12.5, 0.45, size=18, bold=True, color=SLATE)
    if sub:
        txt(slide, sub, 0.35, 1.12, 12.5, 0.3, size=9, color=GRAY)

def bar(slide, x, y, w, label, value, pct, col):
    txt(slide, label, x, y, 2.2, 0.28, size=9.5, color=SLATE)
    rect(slide, x+2.3, y+0.05, w, 0.2, fill=RGBColor(0xE2,0xE8,0xF0))
    rect(slide, x+2.3, y+0.05, (pct/100)*w, 0.2, fill=col)
    txt(slide, value, x+2.3+w+0.08, y, 0.6, 0.28, size=9.5, bold=True, color=SLATE)

# ══════════════════════════════════════════════
# S1 — COVER
# ══════════════════════════════════════════════
s1 = prs.slides.add_slide(BLANK)
rect(s1, 0, 0, 13.33, 7.5, fill=NAVY)
rect(s1, 0, 7.3, 13.33, 0.2, fill=PURPLE)
rect(s1, 0.4, 0.35, 8.0, 0.55, fill=RGBColor(0x2D,0x2B,0x5E))
txt(s1, "प्राण", 0.5, 0.38, 1.2, 0.45, size=14, color=LAVEND)
txt(s1, "Prāṇa — Sanskrit for life force.   PRANA is essential to our career.", 1.9, 0.42, 6.3, 0.38, size=9, color=RGBColor(0xC4,0xB5,0xFD))
txt(s1, "Every worker deserves proof of their career.", 0.4, 1.05, 10, 0.7, size=30, bold=True, color=WHITE)
rect(s1, 0.4, 1.82, 0.6, 0.05, fill=PURPLE)
txt(s1, "→  Employer-Pushed  →  AI-Processed  →  Employee-Owned  career document vault", 0.4, 2.0, 11, 0.35, size=11, color=LAVEND)
txt(s1, "   Global Alumni Portal — your vault survives your resignation. Always on. Always yours.", 0.4, 2.45, 11, 0.35, size=11, color=WHITE)
txt(s1, "Built for   Trust  ·  Compliance  ·  Permanence", 0.4, 2.9, 8, 0.35, size=11, color=LAVEND)
for i,(letter, rest) in enumerate([("P","ersonal"),("R","epository and"),("A","uthenticated"),("N","etworked"),("A","uthentications")]):
    y = 3.55 + i*0.45
    txt(s1, letter, 0.4, y, 0.25, 0.4, size=16, bold=True, color=LAVEND)
    txt(s1, rest, 0.65, y+0.05, 3, 0.35, size=10, color=RGBColor(0x94,0xA3,0xB8))
txt(s1, "Prepared by: Nilesh Shah     ·     Classification: Private & Confidential", 0.4, 5.9, 10, 0.3, size=8.5, color=RGBColor(0x94,0xA3,0xB8))
txt(s1, "PRANA", 11.8, 0.2, 1.3, 0.35, size=13, bold=True, color=WHITE)

# ══════════════════════════════════════════════
# S2 — TOC
# ══════════════════════════════════════════════
s2 = prs.slides.add_slide(BLANK)
rect(s2, 0, 0, 13.33, 7.5, fill=LGRAY)
logo(s2)
rect(s2, 0, 0.45, 13.33, 0.9, fill=INDL)
txt(s2, "What we'll cover today", 0.35, 0.52, 10, 0.4, size=18, bold=True, color=INDIGO)
txt(s2, "15 slides · ~20 minutes", 0.35, 0.88, 4, 0.25, size=8.5, color=PURPLE)
toc = [
    ("03","The Problem","₹42,000 Cr lost to document gaps"),
    ("04","Platform Architecture","Employer → AI → Employee"),
    ("05","Alumni Portal","Always-on career identity"),
    ("06","For the CHRO","Vault health · digest · compliance"),
    ("07","For the CFO","Cost pulse · anomaly ack · coverage"),
    ("08","For the CISO","Access log · anomaly feed · zero-trust"),
    ("09","For the Employee","Mobile vault · share · permanence"),
    ("10","AI Pipeline + Guardrails","Local LLM · privacy filter · 6 stages"),
    ("11","Labour Law Compliance","4 Codes · 29 Acts · wages to gratuity"),
    ("12","DPDP Act 2023","5 rights · SLAs · consent model"),
    ("13","Security + Scale","Encryption · KMS · 1 crore employees"),
    ("14","Immutable Audit + Trust","7-year chain · watermark · zero access"),
    ("15","Market + Why Now","Where PRANA fits · 3 forces · pilot offer"),
]
for ci,col in enumerate([toc[:7], toc[7:]]):
    cx = 0.35 + ci*6.55
    for ri,(num,title,sub) in enumerate(col):
        y = 1.5 + ri*0.77
        rect(s2, cx, y, 0.32, 0.32, fill=PURPLE)
        txt(s2, num, cx+0.04, y+0.04, 0.28, 0.28, size=9, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        txt(s2, title, cx+0.42, y, 5.5, 0.25, size=11, bold=True, color=SLATE)
        txt(s2, sub,   cx+0.42, y+0.24, 5.5, 0.22, size=8.5, color=GRAY)
footer(s2, 2)

# ══════════════════════════════════════════════
# S3 — PROBLEM
# ══════════════════════════════════════════════
s3 = prs.slides.add_slide(BLANK)
rect(s3, 0, 0, 13.33, 7.5, fill=LGRAY)
logo(s3)
hdr(s3, "THE PROBLEM", "India's workforce runs on documents no one can find.", "Three roles. Three pain points. One broken system.")
cards3 = [
    ("Worker", '"My old employer won\'t release my documents.\n3 years of career — gone."', "No portability", BGRED, RED),
    ("CHRO",   '"We run SAP, Darwinbox, 2 local HRMS.\nZero single view. Audit season is a nightmare."', "No visibility", BGWARM, AMBER),
    ("CISO",   '"Who accessed which salary slip, from which IP,\nat 2 AM? Zero audit trail."', "No trail", BGCYA, CYAN),
]
for i,(role,quote,badge,bg,ac) in enumerate(cards3):
    x = 0.3 + i*4.35
    rect(s3, x, 1.75, 4.1, 2.75, fill=bg)
    txt(s3, role, x+0.15, 1.85, 3.8, 0.35, size=12, bold=True, color=ac)
    txt(s3, quote, x+0.15, 2.25, 3.8, 1.6, size=9.5, color=SLATE)
    rect(s3, x+0.15, 4.1, 1.3, 0.28, fill=ac)
    txt(s3, badge, x+0.15, 4.12, 1.3, 0.24, size=8.5, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
rect(s3, 0.3, 4.65, 5.9, 1.45, fill=BGRED)
txt(s3, "Rs.42,000 Cr", 0.5, 4.72, 3.5, 0.55, size=28, bold=True, color=RED)
txt(s3, "estimated annual loss — document fraud, re-issuance\ndelays, and compliance penalties across Indian enterprises", 0.5, 5.3, 5.5, 0.65, size=9.5, color=SLATE)
rect(s3, 6.55, 4.65, 6.4, 1.45, fill=BGVIO)
txt(s3, "The Alumni Blind Spot", 6.75, 4.72, 6, 0.3, size=11, bold=True, color=VIOLET)
txt(s3, "15 crore gig & contract workers have no portable document trail when they move between jobs. PRANA solves this permanently — documents survive resignation.", 6.75, 5.08, 6, 0.88, size=9.5, color=SLATE)
footer(s3, 3)

# ══════════════════════════════════════════════
# S4 — ARCHITECTURE
# ══════════════════════════════════════════════
s4 = prs.slides.add_slide(BLANK)
rect(s4, 0, 0, 13.33, 7.5, fill=LGRAY)
logo(s4)
hdr(s4, "PLATFORM ARCHITECTURE", "Employer pushes. AI processes. Employee owns.")
rect(s4, 0.3, 1.75, 3.2, 3.0, fill=INDL)
txt(s4, "Employer", 0.5, 1.85, 2.8, 0.4, size=13, bold=True, color=INDIGO)
txt(s4, "Portal Upload\nHRMS API Push\nSAP / Darwinbox / Keka\nBatch ZIP + CSV", 0.5, 2.35, 2.8, 2.2, size=10, color=INDIGO)
rect(s4, 3.75, 1.75, 5.6, 3.0, fill=DK)
txt(s4, "PRANA VAULT ENGINE", 4.1, 1.9, 5, 0.3, size=11, bold=True, color=LAVEND, align=PP_ALIGN.CENTER)
txt(s4, "Kafka  ·  Temporal  ·  YugabyteDB  ·  AWS KMS  ·  Redis CRDT", 3.95, 2.28, 5.2, 0.28, size=8.5, color=RGBColor(0x94,0xA3,0xB8), align=PP_ALIGN.CENTER)
txt(s4, "Ingest  →  OCR  →  LLM Extract  →  Resolve  →  Insights  →  Route", 3.95, 2.65, 5.2, 0.28, size=9, color=LAVEND, align=PP_ALIGN.CENTER)
txt(s4, "Validate → S3 put → 1 DB write → 1 Kafka publish → 202\nAll processing async. HTTP handler returns in milliseconds.", 3.95, 3.08, 5.2, 0.5, size=8.5, color=RGBColor(0x94,0xA3,0xB8), align=PP_ALIGN.CENTER)
rect(s4, 9.6, 1.75, 3.4, 3.0, fill=LGREEN)
txt(s4, "Employee", 9.8, 1.85, 3.0, 0.4, size=13, bold=True, color=GREEN)
txt(s4, "Mobile App\nAI Insights\nWatermarked Share\nAlumni Access", 9.8, 2.35, 3.0, 2.2, size=10, color=GREEN)
txt(s4, "->", 3.42, 2.95, 0.4, 0.4, size=20, color=GRAY, align=PP_ALIGN.CENTER)
txt(s4, "->", 9.3, 2.95, 0.4, 0.4, size=20, color=GRAY, align=PP_ALIGN.CENTER)
hl4 = [
    ("Async 202 Response", "HTTP handler returns in ms. All processing via Kafka + Temporal — async, durable, retryable.", INDL),
    ("Envelope Encryption", "Every doc encrypted with employee DEK, wrapped by tenant KEK via AWS KMS. Zero plaintext — ever.", LGREEN),
    ("Dual-Region India", "YugabyteDB active-active: Mumbai + Hyderabad. Sub-10ms cross-region sync. 99.99% uptime SLA.", BGCYA),
]
for i,(t,b,bg) in enumerate(hl4):
    x = 0.3 + i*4.35
    rect(s4, x, 5.0, 4.1, 1.5, fill=bg)
    txt(s4, t, x+0.15, 5.08, 3.8, 0.3, size=10, bold=True, color=INDIGO)
    txt(s4, b, x+0.15, 5.42, 3.8, 0.95, size=9, color=SLATE)
footer(s4, 4)

# ══════════════════════════════════════════════
# S5 — ALUMNI
# ══════════════════════════════════════════════
s5 = prs.slides.add_slide(BLANK)
rect(s5, 0, 0, 13.33, 7.5, fill=LGRAY)
logo(s5)
hdr(s5, "ALUMNI PORTAL", "Your career vault survives your resignation.", "Always-on. Always yours. No HR intervention required.")
items5 = [
    ("Post-Exit Access", "Ex-employee logs in with phone OTP. Views salary slips, Form 16, offer letters, relieving letter — anytime, permanently."),
    ("38-Second Retrieval", "Average alumni document retrieval: 38 seconds. Zero calls to HR. Zero email chains. The 2 AM salary slip request — solved."),
    ("Portable Across Employers", "Career spanning Vertex → Indigo → Bluestar. One vault. All employers' docs co-exist, separated by tenant, unified for the employee."),
    ("Employer Control Preserved", "Alumni access is read-only and watermarked. Employer sets retention window. Documents are never editable by the employee."),
]
for i,(t,b) in enumerate(items5):
    y = 1.75 + i*1.2
    rect(s5, 0.3, y, 7.5, 1.05, fill=WHITE)
    txt(s5, t, 0.5, y+0.08, 7.0, 0.3, size=10.5, bold=True, color=INDIGO)
    txt(s5, b, 0.5, y+0.42, 7.2, 0.55, size=9.5, color=SLATE)
rect(s5, 8.15, 1.75, 4.8, 2.4, fill=DK)
txt(s5, "Alumni self-served this week", 8.35, 1.9, 4.4, 0.3, size=8.5, color=RGBColor(0x94,0xA3,0xB8), align=PP_ALIGN.CENTER)
txt(s5, "7", 8.35, 2.15, 4.4, 0.7, size=42, bold=True, color=LAVEND, align=PP_ALIGN.CENTER)
txt(s5, "0 calls to HR  ·  0 emails sent", 8.35, 2.9, 4.4, 0.3, size=8.5, color=RGBColor(0x94,0xA3,0xB8), align=PP_ALIGN.CENTER)
txt(s5, "Docs from 3 previous employers accessed by 1 employee in 90 seconds.", 8.35, 3.3, 4.4, 0.55, size=9, color=WHITE, align=PP_ALIGN.CENTER)
rect(s5, 8.15, 4.35, 4.8, 1.45, fill=DK)
txt(s5, "90+  system-generated notifications", 8.35, 4.48, 4.4, 0.3, size=10.5, bold=True, color=LAVEND, align=PP_ALIGN.CENTER)
txt(s5, "Email  ·  SMS  ·  WhatsApp  ·  Push  ·  Portal Bell  ·  Incident\nZero manual follow-ups", 8.35, 4.85, 4.4, 0.7, size=8.5, color=RGBColor(0xC4,0xB5,0xFD), align=PP_ALIGN.CENTER)
footer(s5, 5)

# ══════════════════════════════════════════════
# S6 — CHRO
# ══════════════════════════════════════════════
s6 = prs.slides.add_slide(BLANK)
rect(s6, 0, 0, 13.33, 7.5, fill=LGRAY)
logo(s6)
hdr(s6, "FOR THE CHRO", "Complete workforce visibility. In 60 seconds.")
rect(s6, 0.3, 1.75, 7.5, 5.3, fill=WHITE)
txt(s6, "CHRO Weekly Vault Snapshot  ·  Week 32 · Aug 2025", 0.5, 1.85, 7, 0.3, size=10, bold=True, color=SLATE)
for i,(v,l,d,dc) in enumerate([("214","Docs processed","+ 18 vs last week",GREEN),("87%","Vault completeness","+ 1.2% this week",GREEN),("4","Exceptions open","2 past SLA",RED),("7","Alumni self-served","0 HR calls",GREEN)]):
    x = 0.5 + i*1.82
    rect(s6, x, 2.25, 1.7, 0.95, fill=LGRAY)
    txt(s6, v, x+0.08, 2.3, 1.55, 0.38, size=18, bold=True, color=SLATE)
    txt(s6, l, x+0.08, 2.62, 1.55, 0.25, size=7.5, color=GRAY)
    txt(s6, d, x+0.08, 2.85, 1.55, 0.22, size=7.5, color=dc)
txt(s6, "Documents by type this week", 0.5, 3.35, 6, 0.25, size=8.5, bold=True, color=GRAY)
for i,(label,cnt,pct,col) in enumerate([("Salary slips","158",74,PURPLE),("Offer letters","38",18,CYAN),("Form 16","13",6,GREEN)]):
    y = 3.65 + i*0.52
    bar(s6, 0.5, y, 4.5, label, cnt, pct, col)
rect(s6, 0.3, 5.25, 7.5, 0.5, fill=BGWARM)
txt(s6, "Warning:  4 exceptions unresolved — 2 past 4-hr SLA.  Oldest: 31 hrs ago", 0.5, 5.33, 7, 0.3, size=9, color=RGBColor(0x92,0x40,0x0E))
rect(s6, 0.3, 5.82, 7.5, 0.5, fill=LGREEN)
txt(s6, "OK:  7 ex-employees downloaded own docs. Zero HR needed. Avg retrieval: 38 seconds.", 0.5, 5.9, 7, 0.3, size=9, color=RGBColor(0x06,0x5F,0x46))
for i,(t,b) in enumerate([("Vault Health by Department","Engineering 86%  ·  Finance 91%  ·  HR 94%  ·  Sales 74%. Know gaps before audit."),("Compliance Calendar","Labour law obligations by deadline. PF return, ESI filing, gratuity nominations — RAG status."),("Weekly / Monthly / Quarterly Digest","Automated digest to CHRO inbox. Every number backed by an immutable event log.")]):
    y = 1.75 + i*1.77
    rect(s6, 8.1, y, 4.9, 1.6, fill=INDL)
    txt(s6, t, 8.3, y+0.1, 4.6, 0.3, size=10, bold=True, color=INDIGO)
    txt(s6, b, 8.3, y+0.48, 4.6, 1.0, size=9.5, color=SLATE)
footer(s6, 6)

# ══════════════════════════════════════════════
# S7 — CFO
# ══════════════════════════════════════════════
s7 = prs.slides.add_slide(BLANK)
rect(s7, 0, 0, 13.33, 7.5, fill=LGRAY)
logo(s7)
hdr(s7, "FOR THE CFO", "Financial document risk — one number, always current.")
rect(s7, 0.3, 1.75, 7.5, 5.3, fill=WHITE)
txt(s7, "CFO Weekly Cost & Compliance Pulse  ·  Week 32", 0.5, 1.85, 7, 0.3, size=10, bold=True, color=SLATE)
for i,(v,l,d,dc) in enumerate([("1,997","Headcount","Budget 2,050 · -53",RED),("11","Joiners","11 docs auto-pushed",GREEN),("6","Exits","Rs.9L replacement",RED),("2","Anomalies","Awaiting CFO ack",AMBER)]):
    x = 0.5 + i*1.82
    rect(s7, x, 2.25, 1.7, 0.95, fill=LGRAY)
    txt(s7, v, x+0.08, 2.3, 1.55, 0.38, size=16, bold=True, color=SLATE)
    txt(s7, l, x+0.08, 2.62, 1.55, 0.25, size=7.5, color=GRAY)
    txt(s7, d, x+0.08, 2.85, 1.55, 0.22, size=7.5, color=dc)
txt(s7, "Financial doc compliance", 0.5, 3.35, 6, 0.25, size=8.5, bold=True, color=GRAY)
for i,(label,cnt,pct,col) in enumerate([("Salary slips","98%",98,CYAN),("Form 16","96%",96,CYAN),("PF statements","81%",81,AMBER)]):
    y = 3.65 + i*0.52
    bar(s7, 0.5, y, 4.5, label, cnt, pct, col)
rect(s7, 0.3, 5.3, 7.5, 0.6, fill=BGWARM)
txt(s7, "Warning:  2 payroll anomalies — 47 employees, Rs.2.3L delta. Awaiting CFO acknowledgement. SLA: 48 hrs.", 0.5, 5.38, 7, 0.42, size=9, color=RGBColor(0x92,0x40,0x0E))
for i,(t,b,bg) in enumerate([("Anomaly Detection + Ack Workflow","Bulk salary revisions, PAN duplicates — flagged automatically. CFO one-click ack with full Temporal audit trail.",BGCYA),("Form 16 Delivery Tracking","1,921 of 2,000 (96%). Remaining 79 are mid-year joiners — not yet due. CFO sees it in digest, not in an audit.",BGCYA),("Rs.33L Monthly Attrition Cost","Surfaced automatically in CFO monthly digest. No analyst required. Backed by immutable workflow log.",BGCYA)]):
    y = 1.75 + i*1.77
    rect(s7, 8.1, y, 4.9, 1.6, fill=bg)
    txt(s7, t, 8.3, y+0.1, 4.6, 0.3, size=10, bold=True, color=CYAN)
    txt(s7, b, 8.3, y+0.48, 4.6, 1.0, size=9.5, color=SLATE)
footer(s7, 7)

# ══════════════════════════════════════════════
# S8 — CISO
# ══════════════════════════════════════════════
s8 = prs.slides.add_slide(BLANK)
rect(s8, 0, 0, 13.33, 7.5, fill=LGRAY)
logo(s8)
hdr(s8, "FOR THE CISO", "Every access. Every device. Every second. Logged.")
rect(s8, 0.3, 1.75, 7.5, 5.3, fill=WHITE)
txt(s8, "InfoSec Weekly Security Pulse  ·  Week 32", 0.5, 1.85, 7, 0.3, size=10, bold=True, color=SLATE)
for i,(v,l,d,dc) in enumerate([("1,847","Doc accesses","100% watermarked",GREEN),("3","Anomalies","2 open",RED),("1","Force-logout","Bulk · Rajkot IP",RED),("34","Share tokens","31 expired",GREEN)]):
    x = 0.5 + i*1.82
    rect(s8, x, 2.25, 1.7, 0.95, fill=LGRAY)
    txt(s8, v, x+0.08, 2.3, 1.55, 0.38, size=16, bold=True, color=SLATE)
    txt(s8, l, x+0.08, 2.62, 1.55, 0.25, size=7.5, color=GRAY)
    txt(s8, d, x+0.08, 2.85, 1.55, 0.22, size=7.5, color=dc)
txt(s8, "Anomaly breakdown", 0.5, 3.35, 6, 0.25, size=8.5, bold=True, color=GRAY)
for i,(label,cnt,pct,col) in enumerate([("Bulk access (HIGH)","36%",36,RED),("Foreign IP (MED)","27%",27,AMBER),("Failed TOTP (LOW)","27%",27,GREEN)]):
    y = 3.65 + i*0.52
    txt(s8, label, 0.5, y, 2.5, 0.28, size=9, color=SLATE)
    rect(s8, 3.1, y+0.05, 4.0, 0.2, fill=RGBColor(0xE2,0xE8,0xF0))
    rect(s8, 3.1, y+0.05, (pct/100)*4.0, 0.2, fill=col)
    txt(s8, cnt, 7.15, y, 0.5, 0.28, size=9.5, bold=True, color=col)
rect(s8, 0.3, 5.28, 7.5, 0.65, fill=BGRED)
txt(s8, "HIGH:  Bulk access — 47 docs in 4 min, Rajkot IP 103.21.x.x. Force-logout applied. Watermark trail preserved.", 0.5, 5.36, 7, 0.45, size=9, color=RGBColor(0x9F,0x12,0x39))
for i,(t,b) in enumerate([("Foreign IP Detection","OA-Operator logged in from Singapore. Flagged immediately. CISO sees full IP. Employee sees city only — by design."),("4.2-Minute Response Time","Detection to force-logout. Avg across Q2: 4.2 minutes. Watermark trail preserved on every access."),("Account Lock Policy","5 failed TOTP = auto-lock. CISO sees: \"Locked — BULK_ACCESS_ANOMALY. Auto-unlocks in 18h 23m.\" Manual override needs audit reason.")]):
    y = 1.75 + i*1.77
    rect(s8, 8.1, y, 4.9, 1.6, fill=LGREEN)
    txt(s8, t, 8.3, y+0.1, 4.6, 0.3, size=10, bold=True, color=GREEN)
    txt(s8, b, 8.3, y+0.48, 4.6, 1.0, size=9.5, color=SLATE)
footer(s8, 8)

# ══════════════════════════════════════════════
# S9 — EMPLOYEE
# ══════════════════════════════════════════════
s9 = prs.slides.add_slide(BLANK)
rect(s9, 0, 0, 13.33, 7.5, fill=LGRAY)
logo(s9)
hdr(s9, "FOR THE EMPLOYEE", "Your career. All of it. In your pocket. Forever.")
rect(s9, 0.3, 1.75, 2.8, 5.2, fill=DK)
txt(s9, "PRANA VAULT", 0.5, 1.9, 2.4, 0.28, size=8.5, bold=True, color=LAVEND, align=PP_ALIGN.CENTER)
txt(s9, "Arjun Kapoor  ·  3 employers", 0.5, 2.18, 2.4, 0.24, size=8, color=GRAY, align=PP_ALIGN.CENTER)
for i,(label,count) in enumerate([("Salary Slips","24"),("Form 16","6"),("Offer Letters","3"),("Relieving Ltr","2")]):
    y = 2.55 + i*0.55
    rect(s9, 0.45, y, 2.5, 0.45, fill=RGBColor(0x2D,0x2B,0x5E))
    txt(s9, label, 0.58, y+0.08, 1.5, 0.28, size=9, color=RGBColor(0xC4,0xB5,0xFD))
    txt(s9, count, 2.18, y+0.08, 0.6, 0.28, size=11, bold=True, color=LAVEND)
rect(s9, 0.45, 4.82, 2.5, 0.75, fill=INDIGO)
txt(s9, "AI Insight", 0.58, 4.88, 2.3, 0.25, size=8.5, bold=True, color=LAVEND)
txt(s9, '"Consistent progression. Senior band for your experience."', 0.58, 5.1, 2.3, 0.4, size=7.5, color=WHITE, italic=True)
txt(s9, "Salary not shown  ·  PAN never displayed", 0.45, 5.68, 2.5, 0.22, size=7, color=GRAY, align=PP_ALIGN.CENTER)
features9 = [
    ("Mobile-First Vault", "Face ID / fingerprint re-auth for every document access. Server-side watermark before bytes reach the device."),
    ("One-Time Secure Share", "Share salary slip with bank. 6-digit OTP to recipient. 10-min session. Auto-expires. Watermarked with recipient name + timestamp."),
    ("Privacy by Design", "Raw salary — never shown. PAN — never displayed. AI insight only. Your financials stay yours."),
    ("Alumni Access Forever", "Left 5 years ago? Login with OTP. All your employers' docs — there. 38 seconds. No calls to HR."),
    ("DPDP Rights Built-In", "Request erasure, export data, correct errors, withdraw consent — all from the mobile app. Within statutory SLA."),
]
for i,(t,b) in enumerate(features9):
    row,col = divmod(i,2)
    if i == 4:
        x, y, w = 3.35, 5.55, 9.65
    else:
        x = 3.35 + col*4.9; y = 1.75 + row*1.82; w = 4.6
    rect(s9, x, y, w, 1.6, fill=WHITE if i<4 else INDL)
    txt(s9, t, x+0.15, y+0.12, w-0.2, 0.3, size=10, bold=True, color=INDIGO)
    txt(s9, b, x+0.15, y+0.48, w-0.2, 1.0, size=9.5, color=SLATE)
footer(s9, 9)

# ══════════════════════════════════════════════
# S10 — AI PIPELINE
# ══════════════════════════════════════════════
s10 = prs.slides.add_slide(BLANK)
rect(s10, 0, 0, 13.33, 7.5, fill=LGRAY)
logo(s10)
hdr(s10, "AI PIPELINE + GUARDRAILS", "Local LLM. Full document in. Insights only out.", "Data sovereignty guaranteed — no document ever leaves your infrastructure boundary.")
stages = [
    ("01","Ingest","Validate · S3 · Kafka",LGRAY),
    ("02","OCR","Tesseract\nto Textract",LGRAY),
    ("03","LLM Extract","Qwen 2.5-14B\nFull doc in",INDL),
    ("04","Resolve","Identity · PAN\ndedup · EPFO",LGRAY),
    ("05","Insights","Llama 3.1-8B\nInsights only out",LGREEN),
    ("06","Route to Vault","SSE · Employee\naccess",LGRAY),
]
bw = 1.95
for i,(num,title,body,bg) in enumerate(stages):
    x = 0.3 + i*(bw+0.08)
    rect(s10, x, 1.75, bw, 2.0, fill=bg)
    txt(s10, num, x+0.08, 1.82, 0.5, 0.25, size=8, bold=True, color=PURPLE)
    txt(s10, title, x+0.08, 2.2, bw-0.1, 0.3, size=10, bold=True, color=SLATE, align=PP_ALIGN.CENTER)
    txt(s10, body, x+0.08, 2.58, bw-0.1, 0.9, size=8.5, color=GRAY, align=PP_ALIGN.CENTER)
    if i < 5:
        txt(s10, ">", x+bw+0.01, 2.55, 0.1, 0.3, size=14, color=GRAY)
rect(s10, 0.3, 4.0, 12.7, 1.45, fill=BGRED)
txt(s10, "PRIVACY FILTER — Non-Negotiable", 0.5, 4.08, 6, 0.3, size=11, bold=True, color=RED)
txt(s10, 'Stage 05 output passes a mandatory privacy filter. Any raw salary figure, PAN, or NIK is BLOCKED.\nResponse: "I can share insights about your career but cannot share specific financial figures."\nHardcoded. Not configurable. Not bypassable.', 0.5, 4.42, 12.3, 0.88, size=9.5, color=RGBColor(0x9F,0x12,0x39))
rect(s10, 0.3, 5.62, 6.2, 1.3, fill=INDL)
txt(s10, "Locally Hosted LLM", 0.5, 5.7, 5.8, 0.3, size=10.5, bold=True, color=INDIGO)
txt(s10, "Qwen 2.5-14B + Llama 3.1-8B run on GPU workers inside your infrastructure. No document ever sent to OpenAI, Anthropic, or any external API. HuggingFace today — Ollama / vLLM on-premise when ready.", 0.5, 6.05, 5.8, 0.78, size=9.5, color=SLATE)
rect(s10, 6.8, 5.62, 6.2, 1.3, fill=LGREEN)
txt(s10, "Multilingual — Hindi + English", 7.0, 5.7, 5.8, 0.3, size=10.5, bold=True, color=GREEN)
txt(s10, "BGE-M3 multilingual embeddings. Salary slips in Hindi, offer letters in English, PF statements in Marathi. India's workforce, in India's languages.", 7.0, 6.05, 5.8, 0.78, size=9.5, color=SLATE)
footer(s10, 10)

# ══════════════════════════════════════════════
# S11 — LABOUR LAW
# ══════════════════════════════════════════════
s11 = prs.slides.add_slide(BLANK)
rect(s11, 0, 0, 13.33, 7.5, fill=LGRAY)
logo(s11)
hdr(s11, "LABOUR LAW COMPLIANCE", "29 Acts. 4 Codes. One vault that covers them all.")
codes = [
    ("Code on Wages","Wage slips  ·  Wage register export\nBonus records  ·  3-year retention lock",INDL,INDIGO),
    ("Industrial Relations","Appointment letters  ·  Retrenchment notices\nGrievance records",BGVIO,VIOLET),
    ("Social Security Code","PF statements  ·  ESI cards\nGratuity Form F / I / L / M  ·  Maternity records",LGREEN,GREEN),
    ("OSH Code","Contract labour records\nMigrant worker passbook  ·  Health certificates",BGCYA,CYAN),
]
for i,(title,body,bg,ac) in enumerate(codes):
    x = 0.3 + i*3.26
    rect(s11, x, 1.75, 3.1, 1.82, fill=bg)
    txt(s11, title, x+0.15, 1.85, 2.8, 0.3, size=10.5, bold=True, color=ac)
    txt(s11, body, x+0.15, 2.22, 2.8, 1.22, size=9.5, color=SLATE)
details11 = [
    ("Gratuity — Critical Gap Closed","Form F (nomination)  ·  Form I (application)  ·  Form L / M (payment / rejection).\nPRANA timestamps all four. Disputes now have an immutable paper trail.",BGWARM),
    ("Form 16 — Repository, Not Issuer","PRANA stores employer-pushed Form 16 (Parts A + B). TRACES is authoritative. PRANA is the employee's permanent copy — accessible post-exit, shareable with bank or CA.",LGREEN),
    ("Maternity Records — Sensitive Category","Enhanced privacy flag. Dual-consent required for any share. Medical-adjacent data treated with strictest access controls under PRANA's privacy architecture.",BGVIO),
]
for i,(title,body,bg) in enumerate(details11):
    x = 0.3 + i*4.35
    rect(s11, x, 3.8, 4.1, 2.45, fill=bg)
    txt(s11, title, x+0.15, 3.9, 3.8, 0.3, size=10.5, bold=True, color=SLATE)
    txt(s11, body, x+0.15, 4.28, 3.8, 1.82, size=9.5, color=SLATE)
rect(s11, 0.3, 6.48, 12.7, 0.42, fill=DK)
txt(s11, "In addition to 4 Labour Codes:  IT Act 2000  ·  Income Tax Act (TDS / Form 16)  ·  DPDP Act 2023  —  all covered within the same PRANA vault.", 0.5, 6.54, 12.3, 0.28, size=9, color=LAVEND)
footer(s11, 11)

# ══════════════════════════════════════════════
# S12 — DPDP
# ══════════════════════════════════════════════
s12 = prs.slides.add_slide(BLANK)
rect(s12, 0, 0, 13.33, 7.5, fill=LGRAY)
logo(s12)
hdr(s12, "DPDP ACT 2023", "Employee data rights — built-in, not bolted-on.")
rights_r = [("Erasure","30 days SLA","30d",RED),("Export","72 hour SLA","72h",AMBER),("Correction","15 days SLA","15d",PURPLE),("Consent Withdrawal","Immediate","Instant",GREEN),("Grievance","48h ack · 30d","48h",CYAN)]
for i,(title,sla,badge,col) in enumerate(rights_r):
    x = 0.3 + i*2.56
    rect(s12, x, 1.75, 2.4, 1.62, fill=WHITE)
    txt(s12, title, x+0.1, 1.85, 2.2, 0.3, size=10, bold=True, color=SLATE, align=PP_ALIGN.CENTER)
    txt(s12, sla, x+0.1, 2.2, 2.2, 0.25, size=8.5, color=GRAY, align=PP_ALIGN.CENTER)
    rect(s12, x+0.65, 2.52, 1.1, 0.28, fill=col)
    txt(s12, badge, x+0.65, 2.54, 1.1, 0.24, size=8.5, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
rect(s12, 0.3, 3.6, 6.3, 2.9, fill=DK)
txt(s12, "Consent Model", 0.5, 3.7, 5.8, 0.3, size=11, bold=True, color=LAVEND)
txt(s12, "Consent is per-tenant, per-purpose — not a global toggle.\nEvery grant and withdrawal timestamped in consent_log.\nWithdrawal is immediate — processing stops, Kafka notifies all services within milliseconds.\n\nAudit logs survive erasure requests.\n7-year legal retention overrides DPDP — told to employees transparently.", 0.5, 4.08, 5.8, 2.28, size=9.5, color=WHITE)
for i,(t,b) in enumerate([("100% DPDP Compliance — Q2","All requests within SLA. 1 erasure completed Day 12 of 30. 1 export completed Day 3 of 3. Both audit-logged with full Temporal workflow trail."),("Temporal Workflow — Not a Manual Process","Every right triggers a named workflow: ErasureWorkflow, DataExportWorkflow, ConsentWithdrawalWorkflow. Durable, retryable, auditable. Never a manual ticket.")]):
    y = 3.6 + i*1.48
    rect(s12, 6.85, y, 6.15, 1.32, fill=INDL)
    txt(s12, t, 7.05, y+0.1, 5.8, 0.3, size=10, bold=True, color=INDIGO)
    txt(s12, b, 7.05, y+0.48, 5.8, 0.75, size=9.5, color=SLATE)
footer(s12, 12)

# ══════════════════════════════════════════════
# S13 — SECURITY + SCALE
# ══════════════════════════════════════════════
s13 = prs.slides.add_slide(BLANK)
rect(s13, 0, 0, 13.33, 7.5, fill=LGRAY)
logo(s13)
hdr(s13, "SECURITY + SCALE", "Enterprise-grade encryption. India-scale infrastructure.")
txt(s13, "Three-Key Encryption Model", 0.35, 1.75, 6, 0.3, size=11, bold=True, color=SLATE)
enc = [("Document",RGBColor(0xFE,0xF3,0xC7)),("Employee DEK",INDL),("Tenant KEK",BGVIO),("AWS KMS\nap-south-1",LGREEN)]
for i,(label,col) in enumerate(enc):
    x = 0.35 + i*2.82
    rect(s13, x, 2.1, 2.5, 0.75, fill=col)
    txt(s13, label, x+0.1, 2.2, 2.3, 0.55, size=9.5, bold=True, color=SLATE, align=PP_ALIGN.CENTER)
    if i < 3:
        txt(s13, "->", x+2.55, 2.35, 0.25, 0.3, size=14, color=GRAY)
for i,note in enumerate(["Passwords  —  Argon2id  ·  time=2  ·  memory=65,536  ·  parallelism=2","PAN  —  FF3-1 Format-Preserving Encryption per employee DEK","TOTP  —  AES-256-GCM  ·  Never stored in plaintext","Cache  —  pan_token (HMAC-SHA256) only  ·  No plaintext PAN ever cached"]):
    txt(s13, "·  " + note, 0.35, 3.02+i*0.38, 6.2, 0.3, size=9.5, color=SLATE)
txt(s13, "Built for India's Workforce — All of It", 6.85, 1.75, 6, 0.3, size=11, bold=True, color=SLATE)
for i,(v,l) in enumerate([("1L+","organisations"),("1 Cr+","employees"),("125 PB","document storage"),("53","durable workflows")]):
    row,col = divmod(i,2)
    x = 6.85 + col*2.85; y = 2.1 + row*1.05
    rect(s13, x, y, 2.65, 0.9, fill=WHITE)
    txt(s13, v, x+0.12, y+0.05, 2.4, 0.45, size=20, bold=True, color=INDIGO)
    txt(s13, l, x+0.12, y+0.52, 2.4, 0.28, size=9, color=GRAY)
for i,(t,b) in enumerate([("YugabyteDB","Dual-region Mumbai + Hyderabad  ·  256 tablets  ·  99.99% uptime SLA  ·  PostgreSQL-compatible"),("Apache Kafka","5 topics  ·  12 partitions  ·  AWS MSK KRaft mode  ·  MirrorMaker 2 bidirectional cross-region sync"),("Redis CRDT","ElastiCache Global Datastore  ·  Active-active  ·  Sub-10ms cross-region sync  ·  JWT revocation namespace")]):
    y = 4.5 + i*0.82
    rect(s13, 0.35, y, 12.6, 0.7, fill=WHITE)
    txt(s13, t, 0.55, y+0.08, 3.5, 0.3, size=10, bold=True, color=INDIGO)
    txt(s13, b, 4.2, y+0.12, 8.5, 0.28, size=9.5, color=SLATE)
footer(s13, 13)

# ══════════════════════════════════════════════
# S14 — AUDIT + TRUST
# ══════════════════════════════════════════════
s14 = prs.slides.add_slide(BLANK)
rect(s14, 0, 0, 13.33, 7.5, fill=LGRAY)
logo(s14)
hdr(s14, "IMMUTABLE AUDIT + TRUST", "We cannot see your salary. We cannot access your docs. By design.")
txt(s14, "The Audit Chain — Every Event, Forever", 0.35, 1.75, 6.5, 0.3, size=11, bold=True, color=SLATE)
chain14 = [
    ("Document pushed by employer","tenant_id · document_id · actor · IP · timestamp","sha256:a3f8c2...","Jun 24, 09:14",INDL),
    ("AI pipeline: Stage 03 extraction","Qwen 2.5-14B  ·  structured JSON  ·  no raw figures stored","sha256:b7d1e4...","Jun 24, 09:14",LGREEN),
    ("Employee viewed document","Mobile  ·  Face ID auth  ·  watermark applied  ·  IP logged","sha256:c9f2a1...","Jun 24, 11:02",BGWARM),
    ("Share token issued to bank","Recipient OTP  ·  10-min session  ·  watermark: HDFC + timestamp","sha256:d4e8b3...","Jun 24, 14:37",BGVIO),
]
for i,(title,detail,h,ts,bg) in enumerate(chain14):
    y = 2.1 + i*1.1
    rect(s14, 0.35, y, 6.5, 1.0, fill=bg)
    txt(s14, title, 0.55, y+0.08, 5.5, 0.3, size=10, bold=True, color=SLATE)
    txt(s14, detail, 0.55, y+0.42, 4.8, 0.28, size=8.5, color=GRAY)
    txt(s14, h, 0.55, y+0.7, 2.5, 0.22, size=7.5, color=PURPLE)
    txt(s14, ts, 5.2, y+0.08, 1.5, 0.25, size=8, color=GRAY)
rect(s14, 7.1, 1.75, 5.9, 3.62, fill=DK)
txt(s14, "7-Year Retention. Non-Negotiable.", 7.3, 1.85, 5.5, 0.3, size=11, bold=True, color=LAVEND)
rect(s14, 7.3, 2.25, 2.4, 1.1, fill=INDIGO)
txt(s14, "Hot\n0-2 yrs\nYugabyteDB\nQueryable", 7.45, 2.32, 2.1, 0.9, size=9, color=WHITE, align=PP_ALIGN.CENTER)
txt(s14, "->", 9.75, 2.65, 0.4, 0.3, size=14, color=LAVEND)
rect(s14, 10.2, 2.25, 2.4, 1.1, fill=VIOLET)
txt(s14, "Cold\n2-7 yrs\nIceberg on S3\nTamper-evident", 10.35, 2.32, 2.1, 0.9, size=9, color=WHITE, align=PP_ALIGN.CENTER)
txt(s14, "Audit logs survive erasure requests. Legal retention overrides DPDP — told to employees transparently.", 7.3, 3.45, 5.5, 0.75, size=9.5, color=WHITE)
rect(s14, 7.1, 5.55, 5.9, 0.82, fill=BGRED)
txt(s14, "Zero Access to Your Data", 7.3, 5.62, 5.5, 0.3, size=10.5, bold=True, color=RED)
txt(s14, "LLM output = insights only. Raw salary never in DB. PAN never in logs, cache, or API response. Not policy — it is the architecture.", 7.3, 5.95, 5.5, 0.35, size=9, color=SLATE)
rect(s14, 7.1, 6.42, 5.9, 0.68, fill=LGREEN)
txt(s14, "Watermark on Every Share", 7.3, 6.49, 5.5, 0.28, size=10.5, bold=True, color=GREEN)
txt(s14, "Recipient name + timestamp + token. Applied server-side. Cannot be removed. Exfiltration trail always preserved.", 7.3, 6.76, 5.5, 0.26, size=9, color=SLATE)
footer(s14, 14)

# ══════════════════════════════════════════════
# S15 — MARKET + WHY NOW
# ══════════════════════════════════════════════
s15 = prs.slides.add_slide(BLANK)
rect(s15, 0, 0, 13.33, 7.5, fill=LGRAY)
logo(s15)
hdr(s15, "MARKET + WHY NOW", "The gap is real. The moment is now.")
txt(s15, "Where PRANA fits vs. the alternatives", 0.35, 1.75, 6, 0.28, size=10.5, bold=True, color=SLATE)
comps = [
    ("DigiLocker","Government Docs Only","No employer push  ·  No AI processing  ·  No career insights",RGBColor(0xF1,0xF5,0xF9),GRAY,GRAY),
    ("HRMS (SAP / Darwinbox / Keka)","Employer System of Record","Document storage is a side feature  ·  No portability on exit  ·  No employee-owned vault",BGWARM,AMBER,SLATE),
    ("PRANA","India's first employer-pushed  ·  AI-processed  ·  employee-owned career document vault","Employer pushes, employee owns  ·  AI insights, privacy-first  ·  Alumni access forever",DK,LAVEND,WHITE),
]
for i,(name,tagline,points,bg,tc,bc) in enumerate(comps):
    y = 2.1 + i*1.55
    rect(s15, 0.35, y, 6.2, 1.42, fill=bg)
    txt(s15, name, 0.55, y+0.1, 5.8, 0.3, size=11, bold=True, color=tc)
    txt(s15, tagline, 0.55, y+0.4, 5.8, 0.3, size=8.5, color=bc)
    txt(s15, points, 0.55, y+0.72, 5.8, 0.55, size=8.5, color=bc)
txt(s15, "Three forces making this the right moment", 6.85, 1.75, 6.1, 0.28, size=10.5, bold=True, color=SLATE)
forces = [
    ("DPDP Act 2023 Deadline Pressure","Every employer handling employee data must comply. Fines up to Rs.250 Cr. PRANA is the fastest path — not a project, a SaaS subscription."),
    ("15 Crore Workers With No Document Trail","India's gig and contract workforce — fastest-growing, most document-vulnerable. PRANA is the only portable solution."),
    ("HRMS Fragmentation — Peak Pain","Average enterprise runs 3-4 HR systems with zero interoperability. PRANA ingests from all — no rip-and-replace required."),
]
for i,(title,body) in enumerate(forces):
    y = 2.1 + i*1.35
    rect(s15, 6.85, y, 6.15, 1.22, fill=WHITE)
    txt(s15, title, 7.05, y+0.1, 5.85, 0.3, size=10, bold=True, color=SLATE)
    txt(s15, body, 7.05, y+0.45, 5.85, 0.68, size=9.5, color=GRAY)
rect(s15, 6.85, 6.18, 6.15, 0.85, fill=DK)
txt(s15, "Start with one department. Own your workforce's document story in 30 days.\nNo HRMS replacement. No IT project. One API call.", 7.05, 6.25, 5.85, 0.7, size=9.5, color=WHITE, align=PP_ALIGN.CENTER)
footer(s15, 15)

# ─── SAVE ───
out = r"C:\Nilesh\claude-code\prana-docs\PRANA_CXO_Deck.pptx"
prs.save(out)
print(f"Saved: {out}")
